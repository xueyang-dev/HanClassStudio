from __future__ import annotations

import io
import json
import zipfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

from fastapi.testclient import TestClient
from pptx import Presentation

import hcs_api.main as main
import hcs_api.storage as storage
from hcs_api.main import app
from hcs_api.models import AssetCandidate, AssetFile, AssetManifest, ContentBlock, ImageProviderSettings, LLMProviderSettings, LessonBlueprint, LessonProfile, LessonSlide, ProviderSettings, QualityReport, SlideComponent, SourceMaterial
from hcs_api.providers import ProviderError
from hcs_api.strategist import build_interaction_plan, build_media_plan


def test_root_renders_chinese_console() -> None:
    client = TestClient(app)
    response = client.get("/")
    assert response.status_code == 200
    assert "后端控制台" in response.text
    assert "打开前端工作台" in response.text
    assert "模型与 API 设置" in response.text
    assert "LLM API" in response.text
    assert "文生图模型" in response.text
    assert "文生音频 / TTS" in response.text


def test_health_route() -> None:
    client = TestClient(app)
    response = client.get("/api/health")
    assert response.status_code == 200
    assert response.json() == {"status": "ok"}


def test_upload_rejects_formats_outside_frontend_accept_contract(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", runtime_dir / "projects")
    monkeypatch.setattr(main, "PROJECTS_DIR", runtime_dir / "projects")
    response = TestClient(app).post(
        "/api/projects/upload",
        files={"file": ("lesson.txt", b"not a course source", "text/plain")},
    )
    assert response.status_code == 400
    assert "supported" in response.json()["detail"]


def test_cors_allows_vite_fallback_port() -> None:
    response = TestClient(app).options(
        "/api/health",
        headers={
            "Origin": "http://127.0.0.1:5174",
            "Access-Control-Request-Method": "GET",
        },
    )
    assert response.status_code == 200
    assert response.headers["access-control-allow-origin"] == "http://127.0.0.1:5174"

    extra_dev_port = TestClient(app).options(
        "/api/health",
        headers={
            "Origin": "http://127.0.0.1:5175",
            "Access-Control-Request-Method": "GET",
        },
    )
    assert extra_dev_port.status_code == 200
    assert extra_dev_port.headers["access-control-allow-origin"] == "http://127.0.0.1:5175"


def test_agent_skill_docs_exist() -> None:
    root = Path(__file__).resolve().parents[3]
    assert "skills/hanclassstudio/SKILL.md" in (root / "AGENTS.md").read_text(encoding="utf-8")
    assert (root / "CLAUDE.md").read_text(encoding="utf-8").strip() == "See `AGENTS.md`."
    assert "Strict Pipeline" in (root / "skills" / "hanclassstudio" / "SKILL.md").read_text(encoding="utf-8")
    assert "main-generation" in (root / "skills" / "hanclassstudio" / "workflows" / "routing.md").read_text(encoding="utf-8")


def test_provider_settings_round_trip(tmp_path, monkeypatch) -> None:
    monkeypatch.setattr(storage, "CONFIG_DIR", tmp_path / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", tmp_path / "config" / "provider_settings.json")
    client = TestClient(app)
    payload = {
        "llm": {
            "provider": "openai_compatible",
            "base_url": "https://api.example.com/v1",
            "api_key": "test-key",
            "model": "demo-llm",
        },
        "image": {
            "provider": "comfyui",
            "endpoint_url": "http://127.0.0.1:8188",
            "api_key": "",
            "model": "demo-image-workflow",
        },
        "audio": {
            "provider": "openai_tts",
            "endpoint_url": "https://api.example.com/v1/audio",
            "api_key": "audio-key",
            "model": "demo-tts",
            "voice": "demo-voice",
        },
    }

    save_response = client.put("/api/settings/providers", json=payload)
    assert save_response.status_code == 200
    saved = save_response.json()
    # The response contract must preserve non-secret fields while never
    # returning the submitted credentials.
    for key in ("llm", "image", "audio"):
        assert saved[key]["provider"] == payload[key]["provider"]
        assert saved[key]["api_key_present"] is bool(payload[key]["api_key"])
        assert "api_key" not in saved[key]
    assert "capabilities" in saved and "ocr" in saved and "video" in saved
    assert '"api_key"' not in json.dumps(saved)

    get_response = client.get("/api/settings/providers")
    assert get_response.status_code == 200
    fetched = get_response.json()
    for key in ("llm", "image", "audio"):
        assert fetched[key]["provider"] == payload[key]["provider"]
        assert fetched[key]["api_key_present"] is bool(payload[key]["api_key"])
        assert "api_key" not in fetched[key]
    assert "capabilities" in fetched and "ocr" in fetched and "video" in fetched
    assert '"api_key"' not in json.dumps(fetched)


def test_provider_settings_never_expose_config_file_or_credentials(tmp_path, monkeypatch) -> None:
    monkeypatch.setattr(storage, "CONFIG_DIR", tmp_path / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", tmp_path / "config" / "provider_settings.json")
    client = TestClient(app)
    secret = "audit-secret-never-returned"
    response = client.put(
        "/api/settings/providers",
        json={"llm": {"provider": "openai_compatible", "api_key": secret}},
    )
    assert response.status_code == 200
    assert secret not in response.text
    assert "api_key" not in response.json()["llm"]
    assert response.json()["llm"]["api_key_present"] is True

    fetched = client.get("/api/settings/providers")
    assert fetched.status_code == 200
    assert secret not in fetched.text
    assert '"api_key"' not in json.dumps(fetched.json())
    assert fetched.json()["llm"]["api_key_present"] is True

    # The runtime static mount is limited to project artifacts, never config.
    assert client.get("/runtime/config/provider_settings.json").status_code == 404


def test_provider_capability_contract_marks_unimplemented_provider_unavailable(tmp_path, monkeypatch) -> None:
    monkeypatch.setattr(storage, "CONFIG_DIR", tmp_path / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", tmp_path / "config" / "provider_settings.json")
    client = TestClient(app)
    response = client.get("/api/settings/providers/capabilities")
    assert response.status_code == 200
    descriptors = response.json()
    assert {item["capability"] for item in descriptors} == {"llm", "image", "tts", "ocr", "video"}
    deterministic = next(item for item in descriptors if item["provider_id"] == "deterministic")
    assert deterministic["implemented"] is True
    assert deterministic["configured"] is True
    assert deterministic["available"] is True
    assert deterministic["official_url"] == "https://github.com/xueyang-dev/HanClassStudio"
    assert deterministic["license_name"] == "MIT"
    openai = next(item for item in descriptors if item["provider_id"] == "openai_compatible")
    assert openai["api_signup_url"] == "https://platform.openai.com/"
    assert openai["official_url"].startswith("https://platform.openai.com/")
    paddle = next(item for item in descriptors if item["provider_id"] == "paddle_ocr")
    assert paddle["official_url"] == "https://github.com/PaddlePaddle/PaddleOCR"
    assert paddle["license_name"] == "Apache-2.0"
    runway = next(item for item in descriptors if item["provider_id"] == "runway")
    assert runway["implemented"] is False
    assert runway["configurable"] is False
    assert runway["available"] is False
    assert runway["unavailable_reason"]
    image_ids = {item["provider_id"] for item in descriptors if item["capability"] == "image" and item["configurable"]}
    assert image_ids == {"placeholder", "openai_images", "experimental_openai_images", "codex_image"}


def test_provider_settings_partial_update_preserves_omitted_sections(tmp_path, monkeypatch) -> None:
    monkeypatch.setattr(storage, "CONFIG_DIR", tmp_path / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", tmp_path / "config" / "provider_settings.json")
    client = TestClient(app)
    initial = {
        "llm": {"provider": "openai_compatible", "base_url": "https://llm.test/v1", "api_key": "keep-me", "model": "teacher-model"},
        "image": {"provider": "placeholder", "endpoint_url": "", "api_key": "", "model": "placeholder-svg"},
        "audio": {"provider": "placeholder", "endpoint_url": "", "api_key": "", "model": "placeholder-tone", "voice": "default"},
    }
    assert client.put("/api/settings/providers", json=initial).status_code == 200
    updated = client.put("/api/settings/providers", json={"image": {"provider": "placeholder"}})
    assert updated.status_code == 200
    body = updated.json()
    assert body["llm"]["provider"] == initial["llm"]["provider"]
    assert body["llm"]["model"] == initial["llm"]["model"]
    assert body["llm"]["api_key_present"] is True
    assert body["audio"]["provider"] == initial["audio"]["provider"]
    assert body["audio"]["model"] == initial["audio"]["model"]


def test_provider_settings_atomic_write_survives_concurrent_reads(tmp_path, monkeypatch) -> None:
    monkeypatch.setattr(storage, "CONFIG_DIR", tmp_path / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", tmp_path / "config" / "provider_settings.json")
    settings = ProviderSettings(llm=LLMProviderSettings(provider="deterministic", model="deterministic-v1"))
    storage.write_provider_settings(settings)

    def read_or_write(index: int) -> ProviderSettings:
        if index % 2:
            storage.write_provider_settings(settings)
        return storage.read_provider_settings()

    with ThreadPoolExecutor(max_workers=8) as executor:
        results = list(executor.map(read_or_write, range(24)))

    assert all(result.llm.provider == "deterministic" for result in results)


def test_provider_capability_empty_api_key_does_not_erase_existing_secret(tmp_path, monkeypatch) -> None:
    monkeypatch.setattr(storage, "CONFIG_DIR", tmp_path / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", tmp_path / "config" / "provider_settings.json")
    client = TestClient(app)
    initial = {
        "llm": {"provider": "openai_compatible", "base_url": "https://llm.test/v1", "api_key": "secret", "model": "teacher"},
    }
    assert client.put("/api/settings/providers", json=initial).status_code == 200
    updated = client.put(
        "/api/settings/providers",
        json={"capabilities": {"llm": {"providerId": "openai_compatible", "values": {"api_key": "", "model": "teacher-2"}}}},
    )
    assert updated.status_code == 200
    assert updated.json()["llm"]["api_key_present"] is True
    assert '"api_key"' not in json.dumps(updated.json())


def test_project_pipeline_route_runs_full_generation(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(storage, "CONFIG_DIR", runtime_dir / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", runtime_dir / "config" / "provider_settings.json")
    pptx_path = tmp_path / "lesson.pptx"
    _make_pptx(pptx_path)

    client = TestClient(app)
    with pptx_path.open("rb") as file:
        upload_response = client.post(
            "/api/projects/upload",
            files={"file": ("lesson.pptx", file, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
    assert upload_response.status_code == 200
    project_id = upload_response.json()["project_id"]
    project_root = projects_dir / project_id
    for relative_dir in [
        "uploads",
        "sources",
        "analysis",
        "specs",
        "blueprints",
        "assets/images",
        "assets/audio",
        "assets/video",
        "assets/fonts",
        "assets/data",
        "courseware",
        "quality",
        "exports",
        "agent",
        "backup",
    ]:
        assert (project_root / relative_dir).is_dir()
    assert (project_root / "sources" / "source_material.json").exists()

    pipeline_response = client.post(f"/api/projects/{project_id}/pipeline")

    assert pipeline_response.status_code == 200
    body = pipeline_response.json()
    assert body["status"] == "rendered"
    assert body["route"] == "main-generation"
    assert body["quality_state"] == "warning"
    summary_response = client.get(f"/api/projects/{project_id}/design/summary")
    assert summary_response.status_code == 200
    summary = summary_response.json()
    assert summary["learning_state_plan"] is not None
    assert summary["evidence_plan"] is not None
    assert summary["activity_plan"] is not None
    assert "available_actions" in summary
    assert body["lesson_blueprint"]["slides"]
    assert body["asset_manifest"]["audio"]
    assert body["preview_url"] == f"/runtime/projects/{project_id}/courseware/lesson.html"
    assert body["export_url"] == f"/api/projects/{project_id}/export"
    assert (project_root / "specs" / "lesson_spec.md").exists()
    assert (project_root / "specs" / "spec_lock.json").exists()
    assert (project_root / "blueprints" / "lesson_blueprint.json").exists()
    assert (project_root / "blueprints" / "interaction_plan.json").exists()
    assert (project_root / "blueprints" / "media_plan.json").exists()
    assert (project_root / "courseware" / "lesson.html").exists()
    assert (project_root / "courseware" / "render_manifest.json").exists()
    assert (project_root / "quality" / "quality_report.json").exists()

    rerender_response = client.post(
        f"/api/projects/{project_id}/render?expected_revision={body['project_revision']}"
    )
    assert rerender_response.status_code == 200
    rerendered = rerender_response.json()
    assert rerendered["gate_summary"]["export_allowed"] is True
    assert rerendered["export_url"] == f"/api/projects/{project_id}/export"


def test_project_state_exposes_authoritative_stage_and_gate_contract(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(storage, "CONFIG_DIR", runtime_dir / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", runtime_dir / "config" / "provider_settings.json")
    project_id = "state-contract"
    storage.ensure_project(project_id)
    storage.write_model(project_id, "quality_report.json", QualityReport(state="warning", warnings=["review media"]))

    response = TestClient(app).get(f"/api/projects/{project_id}")

    assert response.status_code == 200
    body = response.json()
    assert body["project_revision"] == 0
    assert body["profile_state"] == "inferred"
    assert {stage["stage_id"] for stage in body["stages"]} == {
        "material", "profile", "design", "presentation", "quality", "delivery"
    }
    assert body["gate_summary"]["quality_report"]["state"] == "warning"
    assert body["gate_summary"]["overall_state"] == "warning"
    assert body["gate_summary"]["export_allowed"] is False
    assert body["gate_summary"]["force_export_allowed"] is False


def test_quality_stage_does_not_advertise_render_without_lesson_artifact(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "stage-action-contract"
    storage.ensure_project(project_id)

    body = TestClient(app).get(f"/api/projects/{project_id}").json()
    quality = next(stage for stage in body["stages"] if stage["stage_id"] == "quality")
    assert quality["state"] == "not_started"
    assert quality["available_actions"] == []
    assert "Blueprint artifact is missing" in quality["blockers"]


def test_project_stage_actions_expose_pipeline_and_handoff_operations(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "stage-action-operations"
    storage.ensure_project(project_id)
    storage.write_model(project_id, "lesson_profile.json", LessonProfile(lesson_title="契约"))
    storage.set_profile_state(project_id, "confirmed")

    body = TestClient(app).get(f"/api/projects/{project_id}").json()
    stages = {stage["stage_id"]: stage for stage in body["stages"]}
    assert "run_pipeline" in stages["design"]["available_actions"]
    assert {"agent_package", "agent_validate"}.issubset(stages["delivery"]["available_actions"])


def test_gate_summary_requires_all_four_gates_and_render_before_export(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "all-gates"
    root = storage.ensure_project(project_id)
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(lesson_title="完整门控", slides=[]))
    (root / "courseware" / "lesson.html").write_text("<html>ready</html>", encoding="utf-8")
    storage.write_model(project_id, "quality_report.json", QualityReport(state="pass"))

    before = TestClient(app).get(f"/api/projects/{project_id}")
    assert before.status_code == 200
    before_body = before.json()
    assert before_body["gate_summary"]["overall_state"] == "not_run"
    assert before_body["gate_summary"]["export_allowed"] is False
    assert before_body["gate_summary"]["force_export_allowed"] is False
    assert before_body["gate_summary"]["evidence_alignment"]["state"] == "not_run"
    assert before_body["gate_summary"]["presentation_readiness"]["state"] == "not_run"
    assert before_body["gate_summary"]["presentation_binding"]["state"] == "not_run"

    for relative in (
        "quality/evidence_alignment_report.json",
        "quality/presentation_readiness_report.json",
        "presentation/binding_quality_report.json",
    ):
        storage.write_json(project_id, relative, {"state": "pass"})

    after = TestClient(app).get(f"/api/projects/{project_id}")
    assert after.status_code == 200
    after_body = after.json()
    assert after_body["gate_summary"]["overall_state"] == "passed"
    assert after_body["gate_summary"]["export_allowed"] is True
    assert TestClient(app).get(f"/api/projects/{project_id}/export").status_code == 200


def test_presentation_stage_stays_blocked_when_readiness_gate_blocks(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "readiness-stage-blocked"
    root = storage.ensure_project(project_id)
    storage.write_model(project_id, "source_material.json", SourceMaterial(source_type="pptx", original_filename="lesson.pptx"))
    storage.write_model(project_id, "lesson_profile.json", LessonProfile(lesson_title="课程"))
    storage.set_profile_state(project_id, "confirmed")
    for relative in ("learning/learning_state_plan.json", "learning/evidence_plan.json", "learning/activity_plan.json"):
        storage.write_json(project_id, relative, {"schema": "test"})
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(lesson_title="课程", slides=[]))
    (root / "courseware" / "lesson.html").write_text("<html>ready</html>", encoding="utf-8")
    storage.write_json(project_id, "quality/evidence_alignment_report.json", {"state": "pass"})
    storage.write_json(project_id, "quality/presentation_readiness_report.json", {"state": "blocked", "blocking_reasons": ["missing activity evidence"]})
    storage.write_json(project_id, "presentation/binding_quality_report.json", {"state": "pass"})
    storage.write_model(project_id, "quality_report.json", QualityReport(state="pass"))
    storage.bump_project_revision(project_id)

    body = TestClient(app).get(f"/api/projects/{project_id}").json()
    presentation = next(stage for stage in body["stages"] if stage["stage_id"] == "presentation")
    assert presentation["state"] == "blocked"
    assert "missing activity evidence" in presentation["blockers"]
    assert body["current_stage"] == "presentation"


def test_force_export_cannot_bypass_missing_blueprint_or_render(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "technical-export-blocker"
    storage.ensure_project(project_id)
    storage.write_json(project_id, "quality/evidence_alignment_report.json", {"state": "pass"})
    storage.write_json(project_id, "quality/presentation_readiness_report.json", {"state": "pass"})
    storage.write_json(project_id, "presentation/binding_quality_report.json", {"state": "pass"})
    storage.write_model(project_id, "quality_report.json", QualityReport(state="pass"))

    response = TestClient(app).post(f"/api/projects/{project_id}/export?force=true")

    assert response.status_code == 409
    detail = response.json()["detail"]
    assert detail["code"] == "export_technical_blocked"
    assert detail["force_export_allowed"] is False
    assert any("Blueprint artifact is missing" in reason for reason in detail["blocking_reasons"])


def test_malformed_blueprint_is_a_structured_technical_export_blocker(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "malformed-blueprint"
    root = storage.ensure_project(project_id)
    (root / "blueprints" / "lesson_blueprint.json").write_text("{not-json", encoding="utf-8")
    (root / "courseware" / "lesson.html").write_text("<html>ready</html>", encoding="utf-8")
    for relative in (
        "quality/evidence_alignment_report.json",
        "quality/presentation_readiness_report.json",
        "presentation/binding_quality_report.json",
    ):
        storage.write_json(project_id, relative, {"state": "pass"})
    storage.write_model(project_id, "quality_report.json", QualityReport(state="pass"))

    response = TestClient(app).post(f"/api/projects/{project_id}/export?force=true")

    assert response.status_code == 409
    detail = response.json()["detail"]
    assert detail["code"] == "export_technical_blocked"
    assert any("Blueprint artifact is missing" in reason for reason in detail["blocking_reasons"])


def test_corrupt_render_is_a_non_bypassable_export_blocker(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "corrupt-render"
    root = storage.ensure_project(project_id)
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(lesson_title="损坏渲染", slides=[]))
    (root / "courseware" / "lesson.html").write_bytes(b"\xff\xfe\x00\x01")
    for relative in (
        "quality/evidence_alignment_report.json",
        "quality/presentation_readiness_report.json",
        "presentation/binding_quality_report.json",
    ):
        storage.write_json(project_id, relative, {"state": "pass"})
    storage.write_model(project_id, "quality_report.json", QualityReport(state="pass"))

    response = TestClient(app).post(f"/api/projects/{project_id}/export?force=true")

    assert response.status_code == 409
    detail = response.json()["detail"]
    assert detail["code"] == "export_technical_blocked"
    assert detail["force_export_allowed"] is False
    assert any("unreadable or corrupt" in reason for reason in detail["blocking_reasons"])


def test_stale_quality_report_is_supported_and_blocks_current_preview(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "stale-quality-schema"
    root = storage.ensure_project(project_id)
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(lesson_title="过期质量", slides=[]))
    (root / "courseware" / "lesson.html").write_text("<html>old</html>", encoding="utf-8")
    for relative in (
        "quality/evidence_alignment_report.json",
        "quality/presentation_readiness_report.json",
        "presentation/binding_quality_report.json",
    ):
        storage.write_json(project_id, relative, {"state": "pass"})
    storage.write_model(project_id, "quality_report.json", QualityReport(state="stale", blocking=["profile changed"]))
    storage.bump_project_revision(project_id)

    response = TestClient(app).get(f"/api/projects/{project_id}")

    assert response.status_code == 200
    body = response.json()
    assert body["quality_state"] == "stale"
    assert body["gate_summary"]["quality_report"]["state"] == "stale"
    assert body["gate_summary"]["overall_state"] == "stale"
    assert body["gate_summary"]["export_allowed"] is False
    assert body["preview_url"] is None


def test_export_blocker_is_structured_and_carries_gate_summary(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "structured-blocker"
    root = storage.ensure_project(project_id)
    (root / "courseware" / "lesson.html").write_text("<html></html>", encoding="utf-8")
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(lesson_title="测试", slides=[]))
    storage.write_json(project_id, "quality/evidence_alignment_report.json", {"state": "pass"})
    storage.write_json(project_id, "quality/presentation_readiness_report.json", {"state": "pass"})
    storage.write_json(project_id, "presentation/binding_quality_report.json", {"state": "pass"})
    storage.write_model(project_id, "quality_report.json", QualityReport(state="blocked", blocking=["missing evidence"]))

    response = TestClient(app).get(f"/api/projects/{project_id}/export")

    assert response.status_code == 409
    detail = response.json()["detail"]
    assert detail["code"] == "export_gate_blocked"
    assert "missing evidence" in detail["blocking_reasons"]
    assert detail["gate_summary"]["quality_report"]["state"] == "blocked"
    assert detail["force_export_allowed"] is True


def test_project_listing_and_profile_confirmation_are_persistent(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(storage, "CONFIG_DIR", runtime_dir / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", runtime_dir / "config" / "provider_settings.json")
    project_id = "persistent-profile"
    storage.ensure_project(project_id)
    storage.write_model(project_id, "lesson_profile.json", LessonProfile(lesson_title="初级中文"))
    storage.set_profile_state(project_id, "inferred")
    storage.bump_project_revision(project_id)
    client = TestClient(app)

    before = client.get(f"/api/projects/{project_id}").json()
    confirmed = client.put(
        f"/api/projects/{project_id}/profile",
        json={**before["lesson_profile"], "lesson_title": "确认后的课程"},
    )
    assert confirmed.status_code == 200
    body = confirmed.json()
    assert body["profile_state"] == "confirmed"
    assert body["project_revision"] > before["project_revision"]

    listing = client.get("/api/projects")
    assert listing.status_code == 200
    item = next(entry for entry in listing.json() if entry["project_id"] == project_id)
    assert item["profile_state"] == "confirmed"
    assert item["project_revision"] == body["project_revision"]


def test_project_mutation_rejects_stale_expected_revision(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "revision-conflict"
    storage.ensure_project(project_id)
    storage.write_model(project_id, "lesson_profile.json", LessonProfile(lesson_title="原始课程"))
    storage.set_profile_state(project_id, "inferred")
    storage.bump_project_revision(project_id)
    current = TestClient(app).get(f"/api/projects/{project_id}").json()
    storage.bump_project_revision(project_id)

    response = TestClient(app).put(
        f"/api/projects/{project_id}/profile?expected_revision={current['project_revision']}",
        json={"lesson_title": "不应覆盖"},
    )

    assert response.status_code == 409
    detail = response.json()["detail"]
    assert detail["code"] == "project_revision_conflict"
    assert detail["expected_revision"] == current["project_revision"]
    assert detail["actual_revision"] == current["project_revision"] + 1
    stored = storage.read_model(project_id, "lesson_profile.json", LessonProfile)
    assert stored is not None and stored.lesson_title == "原始课程"


def test_profile_change_invalidates_all_downstream_versions(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "profile-invalidation"
    storage.ensure_project(project_id)
    storage.write_model(project_id, "lesson_profile.json", LessonProfile(lesson_title="旧课程"))
    storage.set_profile_state(project_id, "confirmed")
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(lesson_title="旧课程", slides=[]))
    (projects_dir / project_id / "courseware" / "lesson.html").write_text("<html></html>", encoding="utf-8")
    storage.write_model(project_id, "quality_report.json", QualityReport(state="warning"))
    storage.clear_stale_state(project_id, stages={"profile", "design", "presentation", "media", "render", "quality", "delivery"})

    response = TestClient(app).put(
        f"/api/projects/{project_id}/profile",
        json={"lesson_title": "新课程"},
    )

    assert response.status_code == 200
    body = response.json()
    assert body["profile_state"] == "confirmed"
    assert set(body["stale_state"]["stale_stages"]) == {"design", "presentation", "media", "render", "quality", "delivery"}
    assert body["preview_url"] is None
    assert body["export_url"] is None
    assert body["gate_summary"]["export_allowed"] is False
    stale_render = TestClient(app).post(f"/api/projects/{project_id}/render")
    assert stale_render.status_code == 409
    assert stale_render.json()["detail"]["code"] == "upstream_stale"


def test_persisted_stale_profile_blocks_old_preview_and_export(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "stale-profile-gate"
    root = storage.ensure_project(project_id)
    storage.write_model(project_id, "source_material.json", SourceMaterial(source_type="pptx", original_filename="lesson.pptx"))
    storage.write_model(project_id, "lesson_profile.json", LessonProfile(lesson_title="课程"))
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(lesson_title="课程", slides=[]))
    (root / "courseware" / "lesson.html").write_text("<html>old</html>", encoding="utf-8")
    storage.write_model(project_id, "quality_report.json", QualityReport(state="pass"))
    for relative in (
        "quality/evidence_alignment_report.json",
        "quality/presentation_readiness_report.json",
        "presentation/binding_quality_report.json",
    ):
        storage.write_json(project_id, relative, {"state": "pass"})
    storage.set_profile_state(project_id, "confirmed")
    storage.bump_project_revision(project_id)
    storage.clear_stale_state(project_id, stages={"profile", "design", "presentation", "media", "render", "quality", "delivery"})
    storage.set_profile_state(project_id, "stale")
    storage.clear_stale_state(project_id, stages={"profile", "design", "presentation", "media", "render", "quality", "delivery"})

    body = TestClient(app).get(f"/api/projects/{project_id}").json()

    assert body["profile_state"] == "stale"
    assert body["stale_state"]["stale"] is True
    assert set(body["stale_state"]["stale_stages"]) == {"profile", "design", "presentation", "media", "render", "quality", "delivery"}
    assert body["gate_summary"]["stale"] is True
    assert body["gate_summary"]["export_allowed"] is False
    assert body["preview_url"] is None
    assert body["export_url"] is None


def test_legacy_project_state_is_readable_but_old_outputs_are_not_current(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "legacy-with-outputs"
    root = storage.ensure_project(project_id)
    storage.write_model(project_id, "source_material.json", SourceMaterial(source_type="pptx", original_filename="legacy.pptx"))
    storage.write_model(project_id, "lesson_profile.json", LessonProfile(lesson_title="旧课程"))
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(lesson_title="旧课程", slides=[]))
    (root / "courseware" / "lesson.html").write_text("<html>legacy</html>", encoding="utf-8")
    storage.write_model(project_id, "quality_report.json", QualityReport(state="pass"))
    (root / "exports" / "HanClassStudio_Output_legacy.zip").write_bytes(b"historical")
    for relative in (
        "quality/evidence_alignment_report.json",
        "quality/presentation_readiness_report.json",
        "presentation/binding_quality_report.json",
    ):
        storage.write_json(project_id, relative, {"state": "pass"})
    assert not (root / "assets" / "data" / "project_meta.json").exists()
    assert not (root / "assets" / "data" / "profile_state.json").exists()

    client = TestClient(app)
    response = client.get(f"/api/projects/{project_id}")
    assert response.status_code == 200
    body = response.json()
    # Legacy lineage is intentionally surfaced as stale so the teacher must
    # confirm/re-run before using the historical profile or outputs.
    assert body["profile_state"] == "stale"
    assert body["stale_state"]["stale"] is True
    assert any("lineage is unknown" in reason for reason in body["stale_state"]["reasons"])
    assert body["preview_url"] is None
    assert body["export_url"] is None
    assert body["gate_summary"]["export_allowed"] is False
    assert client.get("/api/projects").status_code == 200
    assert any(item["project_id"] == project_id for item in client.get("/api/projects").json())
    # Reading legacy state must not silently migrate or rewrite source files.
    assert not (root / "assets" / "data" / "project_meta.json").exists()
    assert not (root / "assets" / "data" / "profile_state.json").exists()


def test_media_route_passes_force_regenerate_to_executor(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(storage, "CONFIG_DIR", runtime_dir / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", runtime_dir / "config" / "provider_settings.json")
    project_id = "force-media"
    storage.ensure_project(project_id)
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(lesson_title="媒体", slides=[]))
    captured: dict[str, bool] = {}

    def fake_generate(*_args, **kwargs):
        captured["force_regenerate"] = bool(kwargs.get("force_regenerate"))
        return AssetManifest()

    monkeypatch.setattr(main, "generate_project_media", fake_generate)
    response = TestClient(app).post(f"/api/projects/{project_id}/media?force_regenerate=true")

    assert response.status_code == 200
    assert captured["force_regenerate"] is True


def test_media_review_api_persists_candidate_decision_and_stales_outputs(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "media-review-api"
    root = storage.ensure_project(project_id)
    candidate_path = root / "assets" / "images" / "hero.png"
    candidate_path.write_bytes(b"not-a-real-image")
    candidate = AssetCandidate(id="generated-1", path="assets/images/hero.png", mime_type="image/png", content_hash="hash", source="generated")
    storage.write_model(
        project_id,
        "asset_manifest.json",
        AssetManifest(images=[AssetFile(id="hero", kind="image", path="assets/images/hero.svg", candidates=[candidate], review_state="pending_review")]),
    )

    manifest_response = TestClient(app).get(f"/api/projects/{project_id}/media")
    assert manifest_response.status_code == 200
    decision = TestClient(app).put(
        f"/api/projects/{project_id}/media/hero/review",
        json={"state": "accepted", "candidate_id": "generated-1"},
    )

    assert decision.status_code == 200
    assert decision.json()["review_state"] == "accepted"
    state = TestClient(app).get(f"/api/projects/{project_id}").json()
    assert "render" in state["stale_state"]["stale_stages"]
    assert "quality" in state["stale_state"]["stale_stages"]
    assert "delivery" in state["stale_state"]["stale_stages"]
    quality_stage = next(stage for stage in state["stages"] if stage["stage_id"] == "quality")
    assert "review_media" in quality_stage["available_actions"]
    assert "replace_media" in quality_stage["available_actions"]


def test_unsupported_media_provider_returns_capability_blocker(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(storage, "CONFIG_DIR", runtime_dir / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", runtime_dir / "config" / "provider_settings.json")
    project_id = "unsupported-media"
    storage.ensure_project(project_id)
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(lesson_title="媒体", slides=[]))
    storage.write_provider_settings(ProviderSettings(image=ImageProviderSettings(provider="made_up_provider")))

    response = TestClient(app).post(f"/api/projects/{project_id}/media")

    assert response.status_code == 409
    detail = response.json()["detail"]
    assert detail["code"] == "provider_capability_unavailable"
    assert detail["provider_id"] == "made_up_provider"
    descriptor = next(item for item in TestClient(app).get("/api/settings/providers/capabilities").json() if item["provider_id"] == "made_up_provider")
    assert descriptor["configured"] is False
    assert descriptor["available"] is False


def test_unsupported_llm_provider_does_not_silently_fallback(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(storage, "CONFIG_DIR", runtime_dir / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", runtime_dir / "config" / "provider_settings.json")
    project_id = "unsupported-llm"
    storage.ensure_project(project_id)
    storage.write_model(project_id, "source_material.json", SourceMaterial(source_type="unknown", original_filename="lesson.pdf"))
    storage.write_model(project_id, "lesson_profile.json", LessonProfile())
    storage.write_provider_settings(ProviderSettings(llm=LLMProviderSettings(provider="made_up_llm")))

    response = TestClient(app).post(f"/api/projects/{project_id}/blueprint")

    assert response.status_code == 409
    detail = response.json()["detail"]
    assert detail["code"] == "provider_capability_unavailable"
    assert detail["capability"] == "llm"


def test_unconfigured_llm_provider_returns_structured_capability_blocker(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(storage, "CONFIG_DIR", runtime_dir / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", runtime_dir / "config" / "provider_settings.json")
    project_id = "unconfigured-llm"
    storage.ensure_project(project_id)
    storage.write_model(project_id, "source_material.json", SourceMaterial(source_type="unknown", original_filename="lesson.pdf"))
    storage.write_model(project_id, "lesson_profile.json", LessonProfile())
    storage.write_provider_settings(ProviderSettings(llm=LLMProviderSettings(provider="openai_compatible", base_url="https://llm.test/v1", model="teacher")))

    response = TestClient(app).post(f"/api/projects/{project_id}/blueprint")

    assert response.status_code == 409
    detail = response.json()["detail"]
    assert detail["code"] == "provider_capability_unavailable"
    assert detail["capability"] == "llm"
    assert "credentials" in detail["message"]


def test_provider_execution_failure_does_not_return_success(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(storage, "CONFIG_DIR", runtime_dir / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", runtime_dir / "config" / "provider_settings.json")
    project_id = "provider-execution-failure"
    storage.ensure_project(project_id)
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(lesson_title="媒体", slides=[]))
    storage.write_provider_settings(
        ProviderSettings(image=ImageProviderSettings(provider="openai_images", api_key="configured", model="image")),
    )
    monkeypatch.setattr(main, "generate_project_media", lambda *_args, **_kwargs: (_ for _ in ()).throw(ProviderError("remote unavailable")))

    response = TestClient(app).post(f"/api/projects/{project_id}/media")

    assert response.status_code == 502
    detail = response.json()["detail"]
    assert detail["code"] == "provider_execution_failed"
    assert detail["capability"] == "image"


def test_dependency_invalidation_matrix_matches_pipeline_contract(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    project_id = "invalidation-matrix"
    storage.ensure_project(project_id)
    expected = {
        "ocr": {"profile", "design", "presentation", "media", "render", "quality", "delivery"},
        "profile": {"design", "presentation", "media", "render", "quality", "delivery"},
        "design": {"presentation", "media", "render", "quality", "delivery"},
        "blueprint": {"media", "render", "quality", "delivery"},
        "media": {"render", "quality", "delivery"},
        "render": {"quality", "delivery"},
    }
    for dependency, stages in expected.items():
        storage.clear_stale_state(project_id, stages=set(expected["ocr"]))
        storage.invalidate_downstream(project_id, dependency, dependency)
        payload = storage.read_json(project_id, "assets/data/stale_state.json")
        assert set(payload["stale_stages"]) == stages


def test_golden_sample_pipeline_smoke_exports_expected_artifacts(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(storage, "CONFIG_DIR", runtime_dir / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", runtime_dir / "config" / "provider_settings.json")
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    pptx_path = tmp_path / "hsk1_lesson_14.pptx"
    _make_pptx(pptx_path)

    client = TestClient(app)
    with pptx_path.open("rb") as file:
        upload_response = client.post(
            "/api/projects/upload",
            files={"file": ("hsk1_lesson_14.pptx", file, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
    assert upload_response.status_code == 200
    project_id = upload_response.json()["project_id"]

    pipeline_response = client.post(f"/api/projects/{project_id}/pipeline")
    assert pipeline_response.status_code == 200

    project_root = projects_dir / project_id
    for directory in ["uploads", "sources", "analysis", "learning", "presentation", "specs", "blueprints", "assets", "courseware", "quality", "exports", "agent"]:
        assert (project_root / directory).is_dir()
    for artifact in [
        "learning/learning_state_plan.json",
        "learning/evidence_plan.json",
        "learning/activity_plan.json",
        "presentation/activity_bindings.json",
        "presentation/binding_quality_report.json",
        "specs/lesson_spec.md",
        "specs/spec_lock.json",
        "blueprints/lesson_blueprint.json",
        "blueprints/interaction_plan.json",
        "blueprints/media_plan.json",
        "courseware/lesson.html",
        "quality/quality_report.json",
    ]:
        assert (project_root / artifact).exists()

    quality_report = json.loads((project_root / "quality" / "quality_report.json").read_text(encoding="utf-8"))
    assert "state" in quality_report
    bindings_json = json.loads((project_root / "presentation" / "activity_bindings.json").read_text(encoding="utf-8"))
    binding_report_json = json.loads((project_root / "presentation" / "binding_quality_report.json").read_text(encoding="utf-8"))
    blueprint_json = json.loads((project_root / "blueprints" / "lesson_blueprint.json").read_text(encoding="utf-8"))
    component_ids = [
        component["id"]
        for slide in blueprint_json["slides"]
        for component in slide.get("components", [])
    ]
    assert bindings_json["schema"] == "hanclassstudio.presentation_bindings.v1"
    assert "schema_" not in bindings_json
    assert binding_report_json["schema"] == "hanclassstudio.presentation_bindings.v1"
    assert "schema_" not in binding_report_json
    assert len(component_ids) == len(set(component_ids))
    components = {
        (slide["id"], component["id"])
        for slide in blueprint_json["slides"]
        for component in slide.get("components", [])
    }
    assert all(
        not binding.get("component_id") or (binding["slide_id"], binding["component_id"]) in components
        for binding in bindings_json["bindings"]
    )

    export_path = storage.latest_export_path(project_id)
    assert export_path is not None
    assert export_path.exists()
    with zipfile.ZipFile(export_path) as zf:
        names = set(zf.namelist())
    assert "lesson.html" in names
    assert "assets/data/lesson_blueprint.json" in names
    assert "assets/data/asset_manifest.json" in names
    assert "assets/data/activity_bindings.json" in names
    assert "assets/data/binding_quality_report.json" in names
    assert "assets/data/quality_report.json" in names

    artifacts_response = client.get(f"/api/projects/{project_id}/artifacts")
    assert artifacts_response.status_code == 200
    artifact_body = artifacts_response.json()
    assert artifact_body["spec_lock"]["route"] == "main-generation"
    artifact_paths = {item["path"]: item for group in artifact_body["groups"] for item in group["items"]}
    assert artifact_paths["specs/spec_lock.json"]["exists"] is True
    assert artifact_paths["presentation/activity_bindings.json"]["exists"] is True
    assert artifact_paths["presentation/binding_quality_report.json"]["exists"] is True
    assert artifact_paths["presentation/binding_quality_report.json"]["artifact_type"] == "presentation"
    assert artifact_paths["quality/quality_report.json"]["artifact_type"] == "quality"

    agent_response = client.post(f"/api/projects/{project_id}/agent/package")
    assert agent_response.status_code == 200
    agent_body = agent_response.json()
    assert "AGENTS.md" in agent_body["task_text"]
    assert "courseware/lesson.html" in agent_body["rules_text"]
    assert (project_root / "agent" / "AGENT_TASK.md").exists()
    assert (project_root / "agent" / "AGENT_RULES.md").exists()

    validation_response = client.post(f"/api/projects/{project_id}/agent/validate")
    assert validation_response.status_code == 200
    validation_body = validation_response.json()
    assert validation_body["state"] == "pass"
    assert "All blueprint components are registry-compatible" in validation_body["passed"]
    assert not any("Duplicate component id vocab_cards" in issue for issue in validation_body["blocking"])

    refreshed_artifacts = client.get(f"/api/projects/{project_id}/artifacts").json()
    refreshed_paths = {item["path"]: item for group in refreshed_artifacts["groups"] for item in group["items"]}
    assert refreshed_paths["agent/AGENT_TASK.md"]["exists"] is True


def test_agent_validation_blocks_duplicate_component_ids_with_slide_ids(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "agentdupes"
    project_root = storage.ensure_project(project_id)
    storage.write_json(project_id, "specs/spec_lock.json", {
        "schema": "hanclassstudio.spec_lock.v1",
        "project_id": project_id,
        "route": "main-generation",
        "generation_mode": "guided_redesign",
        "components": {"allowed": ["VocabularyFlipCard"]},
        "quality": {},
    })
    (project_root / "specs" / "lesson_spec.md").write_text("# Lesson", encoding="utf-8")
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(
        lesson_title="重复组件",
        objectives=["x"],
        slides=[
            LessonSlide(id=3, slide_type="VocabularySlide", layout_variant="card_grid", title="你好", components=[
                SlideComponent(id="vocab_cards", component_type="VocabularyFlipCard", data={"items": [{"word": "你好", "pinyin": "nǐ hǎo"}]}),
            ]),
            LessonSlide(id=4, slide_type="VocabularySlide", layout_variant="card_grid", title="您好", components=[
                SlideComponent(id="vocab_cards", component_type="VocabularyFlipCard", data={"items": [{"word": "您好", "pinyin": "nín hǎo"}]}),
            ]),
        ],
    ))
    storage.write_json(project_id, "blueprints/interaction_plan.json", {"schema": "hanclassstudio.interaction_plan.v1", "interactions": []})
    storage.write_json(project_id, "blueprints/media_plan.json", {"schema": "hanclassstudio.media_plan.v1", "images": [], "audio": []})

    validation_response = TestClient(app).post(f"/api/projects/{project_id}/agent/validate")
    assert validation_response.status_code == 200
    validation = validation_response.json()
    assert validation["state"] == "blocked"
    assert any("Duplicate component id vocab_cards" in issue and "3, 4" in issue for issue in validation["blocking"])


def test_editable_pptx_export_after_pipeline_keeps_html_zip_intact(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(storage, "CONFIG_DIR", runtime_dir / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", runtime_dir / "config" / "provider_settings.json")
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    pptx_path = tmp_path / "editable_source.pptx"
    _make_pptx(pptx_path)

    client = TestClient(app)
    with pptx_path.open("rb") as file:
        upload_response = client.post(
            "/api/projects/upload",
            files={"file": ("editable_source.pptx", file, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
    assert upload_response.status_code == 200
    project_id = upload_response.json()["project_id"]
    pipeline_response = client.post(f"/api/projects/{project_id}/pipeline")
    assert pipeline_response.status_code == 200

    pptx_response = client.post(f"/api/projects/{project_id}/export/pptx-editable")
    assert pptx_response.status_code == 200
    body = pptx_response.json()
    assert body["export_type"] == "pptx_editable"
    assert body["editable"] is True
    assert body["interaction_policy"] == "classroom_static_activity"
    assert body["filename"].endswith(".pptx")
    assert body["download_url"].endswith(body["filename"])

    project_root = projects_dir / project_id
    editable_path = project_root / "exports" / body["filename"]
    assert editable_path.exists()
    assert editable_path.stat().st_size > 0
    prs = Presentation(editable_path)
    assert len(prs.slides) >= 1
    assert any(shape.has_text_frame and shape.text.strip() for slide in prs.slides for shape in slide.shapes)
    assert (project_root / "quality" / "pptx_quality_report.json").exists()
    manifest = json.loads((project_root / "exports" / "pptx_export_manifest.json").read_text(encoding="utf-8"))
    assert manifest["forced"] is False

    artifacts_response = client.get(f"/api/projects/{project_id}/artifacts")
    artifact_paths = {item["path"]: item for group in artifacts_response.json()["groups"] for item in group["items"]}
    assert artifact_paths["quality/pptx_quality_report.json"]["exists"] is True
    assert artifact_paths["exports/pptx_export_manifest.json"]["exists"] is True
    assert artifact_paths[f"exports/{body['filename']}"]["exists"] is True

    zip_response = client.get(f"/api/projects/{project_id}/export")
    assert zip_response.status_code == 200
    with zipfile.ZipFile(io.BytesIO(zip_response.content)) as zf:
        names = set(zf.namelist())
    assert "lesson.html" in names
    assert "assets/data/quality_report.json" in names


def test_editable_pptx_export_respects_blocked_quality_and_force(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "blockedpptx"
    storage.ensure_project(project_id)
    storage.write_model(
        project_id,
        "lesson_blueprint.json",
        LessonBlueprint(
            lesson_title="测试课",
            objectives=["完成练习"],
            key_vocabulary=[{"word": "学", "pinyin": "xue2", "meaning": "study"}],
            grammar_points=[],
            slides=[
                LessonSlide(
                    id=1,
                    slide_type="PracticeSlide",
                    layout_variant="basic",
                    title="练习",
                    content_blocks=[ContentBlock(id="c1", text="我学习中文。")],
                    components=[],
                )
            ],
        ),
    )
    (projects_dir / project_id / "courseware" / "lesson.html").write_text("<html>ready</html>", encoding="utf-8")
    storage.write_json(project_id, "quality/evidence_alignment_report.json", {"state": "pass"})
    storage.write_json(project_id, "quality/presentation_readiness_report.json", {"state": "pass"})
    storage.write_json(project_id, "presentation/binding_quality_report.json", {"state": "pass"})
    storage.write_model(project_id, "quality_report.json", QualityReport(state="blocked", blocking=["missing answer"]))

    client = TestClient(app)
    normal_response = client.post(f"/api/projects/{project_id}/export/pptx-editable")
    assert normal_response.status_code == 409

    forced_response = client.post(f"/api/projects/{project_id}/export/pptx-editable?force=true")
    assert forced_response.status_code == 200
    body = forced_response.json()
    manifest = json.loads((projects_dir / project_id / "exports" / "pptx_export_manifest.json").read_text(encoding="utf-8"))
    assert manifest["forced"] is True
    assert "missing answer" in manifest["forced_blockers"]
    assert manifest["force_confirmation"] == "explicit force=true request"
    assert (projects_dir / project_id / "exports" / body["filename"]).exists()


def test_agent_handoff_e2e_validates_then_render_exports(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(storage, "CONFIG_DIR", runtime_dir / "config")
    monkeypatch.setattr(storage, "PROVIDER_SETTINGS_PATH", runtime_dir / "config" / "provider_settings.json")
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    pptx_path = tmp_path / "agent_handoff.pptx"
    _make_pptx(pptx_path)

    client = TestClient(app)
    with pptx_path.open("rb") as file:
        upload_response = client.post(
            "/api/projects/upload",
            files={"file": ("agent_handoff.pptx", file, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
    assert upload_response.status_code == 200
    project_id = upload_response.json()["project_id"]
    project_root = projects_dir / project_id

    blueprint_response = client.post(f"/api/projects/{project_id}/blueprint")
    assert blueprint_response.status_code == 200
    package_response = client.post(f"/api/projects/{project_id}/agent/package")
    assert package_response.status_code == 200
    assert (project_root / "agent" / "AGENT_TASK.md").exists()
    assert (project_root / "agent" / "AGENT_RULES.md").exists()
    assert not (project_root / "courseware" / "lesson.html").exists()
    assert storage.latest_export_path(project_id) is None

    (project_root / "specs" / "lesson_spec.md").write_text(
        (project_root / "specs" / "lesson_spec.md").read_text(encoding="utf-8")
        + "\n\n## Agent Refinement\nAdd a character formation moment for 学 and an audio replay button.",
        encoding="utf-8",
    )
    blueprint = LessonBlueprint.model_validate_json((project_root / "blueprints" / "lesson_blueprint.json").read_text(encoding="utf-8"))
    next_slide_id = len(blueprint.slides) + 1
    blueprint.slides[0].components.append(
        SlideComponent(
            id="agent_audio_replay",
            component_type="AudioButton",
            title="句子复听",
            data={"audio_key": "agent_sentence_audio", "audio_text": "我在学习中文呢。", "label": "播放句子"},
        )
    )
    blueprint.slides.append(
        LessonSlide(
            id=next_slide_id,
            slide_type="CharacterFormationSlide",
            layout_variant="formation",
            title="汉字构形：学",
            content_blocks=[
                ContentBlock(
                    id="agent_char_note",
                    block_type="prompt",
                    text="看部件，想一想“学”的意思。",
                    scaffolding_text="Look at the parts and infer the meaning.",
                )
            ],
            components=[
                SlideComponent(
                    id="agent_character_xue",
                    component_type="CharacterFormation",
                    title="部件到汉字",
                    data={"character": "学", "parts": ["⺍", "冖", "子"], "explanation": "从部件观察汉字结构。"},
                )
            ],
        )
    )
    (project_root / "blueprints" / "lesson_blueprint.json").write_text(blueprint.model_dump_json(indent=2), encoding="utf-8")
    modified_blueprint = LessonBlueprint.model_validate_json((project_root / "blueprints" / "lesson_blueprint.json").read_text(encoding="utf-8"))
    storage.write_json(project_id, "blueprints/interaction_plan.json", build_interaction_plan(modified_blueprint))
    storage.write_json(project_id, "blueprints/media_plan.json", build_media_plan(modified_blueprint))

    validate_response = client.post(f"/api/projects/{project_id}/agent/validate")
    assert validate_response.status_code == 200
    validation = validate_response.json()
    assert validation["state"] == "warning"
    assert "All blueprint components are registry-compatible" in validation["passed"]
    assert any("run render" in item for item in validation["warnings"])
    assert not (project_root / "courseware" / "lesson.html").exists()
    assert storage.latest_export_path(project_id) is None

    render_response = client.post(f"/api/projects/{project_id}/render")
    assert render_response.status_code == 200
    assert (project_root / "courseware" / "lesson.html").exists()
    assert (project_root / "quality" / "quality_report.json").exists()
    # Agent handoff/render does not fabricate the missing State-first gate
    # reports; export remains unavailable until the complete gate contract is
    # run.
    assert storage.latest_export_path(project_id) is None

    export_response = client.get(f"/api/projects/{project_id}/export")
    assert export_response.status_code == 409
    export_detail = export_response.json()["detail"]
    assert export_detail["code"] == "export_gate_blocked"
    assert export_detail["gate_summary"]["overall_state"] == "stale"


def test_blocked_quality_prevents_normal_export_but_force_export_succeeds(tmp_path, monkeypatch) -> None:
    runtime_dir = tmp_path / "runtime"
    projects_dir = runtime_dir / "projects"
    monkeypatch.setattr(storage, "RUNTIME_DIR", runtime_dir)
    monkeypatch.setattr(storage, "PROJECTS_DIR", projects_dir)
    monkeypatch.setattr(main, "PROJECTS_DIR", projects_dir)
    project_id = "blockedproject"
    project_root = storage.ensure_project(project_id)
    (project_root / "courseware" / "lesson.html").write_text("<!doctype html><title>blocked</title>", encoding="utf-8")
    storage.write_model(project_id, "lesson_blueprint.json", LessonBlueprint(lesson_title="阻塞导出", slides=[]))
    storage.write_json(project_id, "quality/evidence_alignment_report.json", {"state": "pass"})
    storage.write_json(project_id, "quality/presentation_readiness_report.json", {"state": "pass"})
    storage.write_json(project_id, "presentation/binding_quality_report.json", {"state": "pass"})
    storage.write_model(
        project_id,
        "quality_report.json",
        QualityReport(state="blocked", blocking=["missing answer"]),
    )

    client = TestClient(app)
    normal_response = client.get(f"/api/projects/{project_id}/export")
    assert normal_response.status_code == 409

    forced_response = client.post(f"/api/projects/{project_id}/export?force=true")
    assert forced_response.status_code == 200
    with zipfile.ZipFile(io.BytesIO(forced_response.content)) as zf:
        names = set(zf.namelist())
        manifest = zf.read("export_manifest.json").decode("utf-8")
    assert "lesson.html" in names
    assert "assets/data/quality_report.json" in names
    assert '"forced": true' in manifest
    assert '"missing answer"' in manifest
    assert '"force_confirmation": "explicit force=true request"' in manifest


def test_component_registry_route_exposes_supported_components() -> None:
    client = TestClient(app)
    response = client.get("/api/component-registry")
    assert response.status_code == 200
    body = response.json()
    assert "VocabularyFlipCard" in body
    assert body["ClassroomGame"]["experimental"] is True


def _make_pptx(path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "第14课 我在学习中文呢"
    slide.placeholders[1].text = "学习\n中文\n我在学习中文呢。"
    prs.save(path)
