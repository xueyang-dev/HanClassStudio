from __future__ import annotations

import io
from pathlib import Path

from PIL import Image
from pptx import Presentation

from hcs_api import storage
from hcs_api.models import (
    AssetFile, AssetManifest, ContentBlock, LessonBlueprint, LessonProfile,
    LessonSlide, MediaRequirements, QualityReport, SlideComponent,
)
from hcs_api.pptx_deck import build_pptx_deck_plan
from hcs_api.pptx_design import MASTER_SOURCE, PROFILE, RECIPES
from hcs_api.pptx_diagnostics import build_contact_sheet
from hcs_api.pptx_exporter import export_editable_pptx


def _png() -> bytes:
    output = io.BytesIO()
    Image.new("RGB", (1600, 900), (224, 154, 112)).save(output, format="PNG")
    return output.getvalue()


def test_master_profile_is_derived_from_real_reference() -> None:
    assert MASTER_SOURCE.endswith("第1课 教学课件 中文 七年级 第一学期.pptx")
    assert (PROFILE.slide_width, PROFILE.slide_height) == (13.333, 7.5)
    assert PROFILE.primary == "5B9BD5"
    assert PROFILE.chinese_hero_size >= 40
    assert PROFILE.minimum_body_size >= 18
    assert {recipe.archetype for recipe in RECIPES.values()} >= {
        "cover", "vocabulary_focus", "formal_informal_contrast",
        "listening_choice", "matching_activity", "recap",
    }


def test_semantic_components_select_activity_archetypes() -> None:
    blueprint = LessonBlueprint(lesson_title="问候", slides=[
        LessonSlide(id=1, slide_type="PracticeSlide", layout_variant="listen_choose", title="听一听", components=[
            SlideComponent(id="listen", component_type="ListenAndChoose", data={"choices": ["你好", "您好"], "answer": "您好", "audio_key": "a"}),
        ]),
        LessonSlide(id=2, slide_type="PracticeSlide", layout_variant="matching", title="连一连", components=[
            SlideComponent(id="match", component_type="MatchGame", data={"pairs": [{"left": "你好", "right": "hello"}]}),
        ]),
        LessonSlide(id=3, slide_type="DialogueSlide", layout_variant="contrast", title="你好 / 您好"),
    ])
    plan = build_pptx_deck_plan(blueprint)
    assert [slide.traditional_layout for slide in plan.slides] == ["listen_choose", "match_pairs", "two_card_contrast"]


def test_master_export_is_editable_legible_and_within_slide(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setattr(storage, "RUNTIME_DIR", tmp_path / "runtime")
    monkeypatch.setattr(storage, "PROJECTS_DIR", tmp_path / "runtime" / "projects")
    project_id = "master_design"
    root = storage.ensure_project(project_id)
    image_path = root / "assets" / "images" / "scene.png"
    image_path.write_bytes(_png())
    blueprint = LessonBlueprint(lesson_title="问候", slides=[
        LessonSlide(id=1, slide_type="CoverSlide", layout_variant="hero", title="老师好！", content_blocks=[
            ContentBlock(id="cover", text="老师好", scaffolding_text="lǎoshī hǎo · Hello, teacher"),
        ], media_requirements=MediaRequirements(image_key="scene", image_prompt="scene")),
        LessonSlide(id=2, slide_type="PracticeSlide", layout_variant="listen_choose", title="听一听，选一选", components=[
            SlideComponent(id="listen", component_type="ListenAndChoose", data={"audio_key": "a", "audio_text": "您好", "choices": ["你好", "您好", "再见"], "answer": "您好"}),
        ]),
    ])
    storage.write_model(project_id, "lesson_blueprint.json", blueprint)
    storage.write_model(project_id, "lesson_profile.json", LessonProfile(lesson_title="问候", learner_level="zero_beginner"))
    storage.write_model(project_id, "asset_manifest.json", AssetManifest(images=[AssetFile(id="scene", kind="image", path="assets/images/scene.png", mime_type="image/png")]))
    storage.write_model(project_id, "quality_report.json", QualityReport(state="pass"))
    path = export_editable_pptx(project_id)
    presentation = Presentation(path)
    assert len(presentation.slides) == 2
    assert any(shape.shape_type == 13 for shape in presentation.slides[0].shapes)
    assert all(
        shape.left >= 0 and shape.top >= 0
        and shape.left + shape.width <= presentation.slide_width
        and shape.top + shape.height <= presentation.slide_height
        for slide in presentation.slides for shape in slide.shapes
    )
    sizes = [run.font.size.pt for slide in presentation.slides for shape in slide.shapes if shape.has_text_frame for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text.strip() and run.font.size]
    assert min(sizes) >= 10
    report = storage.read_json(project_id, "quality/pptx_quality_report.json")
    assert report["off_slide_objects"] == []
    assert report["text_below_minimum"] == []
    assert report["master_profile"]["source"] == MASTER_SOURCE


def test_contact_sheet_is_diagnostic_only(tmp_path: Path) -> None:
    images = []
    for index in range(2):
        path = tmp_path / f"slide-{index}.png"
        Image.new("RGB", (1600, 900), (80 + index * 30, 140, 190)).save(path)
        images.append(path)
    output = build_contact_sheet(images, tmp_path / "diagnostics" / "contact.png")
    assert output.is_file()
    assert Image.open(output).size == (1278, 264)
