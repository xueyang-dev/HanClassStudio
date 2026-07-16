from __future__ import annotations

import json
from datetime import datetime
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from .models import AssetManifest, LessonBlueprint, LessonProfile, QualityReport
from .pptx_design import PROFILE, RECIPES, profile_for_theme
from .presentation_theme import presentation_theme_for_project
from .storage import ensure_project, read_json, read_model, write_json


SUPPORTED_SLIDE_TYPES = {
    "CoverSlide",
    "ObjectiveSlide",
    "VocabularySlide",
    "GrammarPatternSlide",
    "ReadingSlide",
    "DialogueSlide",
    "PracticeSlide",
    "SummarySlide",
    "HomeworkSlide",
}

BG = RGBColor(248, 250, 247)
SURFACE = RGBColor(255, 255, 255)
INK = RGBColor(34, 50, 54)
MUTED = RGBColor(95, 115, 112)
TEAL = RGBColor(8, 126, 139)
MINT = RGBColor(221, 247, 229)
CORAL = RGBColor(242, 95, 92)
GOLD = RGBColor(247, 179, 43)
LINE = RGBColor(220, 232, 226)


import re as _re

PROVIDER_REQUIRED_PPTX = _re.compile(r"provider_required|\[Arabic\]|\[.*?\].*?provider_required")


def _clean_classroom_text(text: str) -> str:
    if PROVIDER_REQUIRED_PPTX.search(text):
        return ""
    return text


def export_editable_pptx(project_id: str, force: bool = False, export_mode: str = "debug") -> Path:
    global PROFILE
    root = ensure_project(project_id)
    theme = presentation_theme_for_project(root)
    # Recipes and coordinates remain reference-master driven; this switches
    # only their shared visual tokens for the current project's export.
    PROFILE = profile_for_theme(theme)
    blueprint = read_model(project_id, "lesson_blueprint.json", LessonBlueprint)
    if not blueprint:
        raise ValueError("Project needs blueprints/lesson_blueprint.json before editable PPTX export")

    report = read_model(project_id, "quality_report.json", QualityReport)
    if not report and not force:
        raise PermissionError("Run quality gate before editable PPTX export")
    if report and report.state == "blocked" and not force:
        raise PermissionError("Quality gate is blocked; pass force=true to export editable PPTX anyway")
    alignment_report = read_json(project_id, "quality/evidence_alignment_report.json") or {}
    if isinstance(alignment_report, dict) and alignment_report.get("state") == "blocked" and not force:
        raise PermissionError("Evidence alignment gate is blocked; pass force=true to export editable PPTX anyway")
    readiness_report = read_json(project_id, "quality/presentation_readiness_report.json") or {}
    if isinstance(readiness_report, dict) and readiness_report.get("state") == "blocked" and not force:
        raise PermissionError("Presentation readiness gate is blocked; pass force=true to export editable PPTX anyway")
    binding_report = read_json(project_id, "presentation/binding_quality_report.json") or {}
    if isinstance(binding_report, dict) and binding_report.get("state") == "blocked" and not force:
        raise PermissionError("Presentation binding gate is blocked; pass force=true to export editable PPTX anyway")

    # Classroom mode: check classroom_quality gate
    is_classroom = export_mode == "classroom"
    if is_classroom:
        from .storage import read_model as _rm
        from .models import ClassroomQualityReport as _CQR
        cqr = _rm(project_id, "classroom_quality_report.json", _CQR)
        if cqr and cqr.state == "blocked" and not force:
            raise PermissionError("Classroom quality gate blocked this export; pass force=true to proceed")

    spec_lock = read_json(project_id, "specs/spec_lock.json") or {}
    interaction_plan = read_json(project_id, "blueprints/interaction_plan.json") or {}
    media_plan = read_json(project_id, "blueprints/media_plan.json") or {}
    manifest = read_model(project_id, "asset_manifest.json", AssetManifest) or AssetManifest()

    # Build PPTX deck plan with kernel + presentation binding artifacts
    from .pptx_deck import build_pptx_deck_plan, build_pptx_structure_report
    from .storage import read_model as _rm, write_json as _wj, read_json as _rj
    from .learner_comprehension import resolve_profile_learner_level
    profile = _rm(project_id, "lesson_profile.json", LessonProfile)
    level = resolve_profile_learner_level(profile) if profile else "zero_beginner"
    # Read kernel artifacts (use read_json for dicts, then construct models)
    ep_data = _rj(project_id, "learning/evidence_plan.json")
    ap_data = _rj(project_id, "learning/activity_plan.json")
    sp_data = _rj(project_id, "learning/learning_state_plan.json")
    binding_data = _rj(project_id, "presentation/activity_bindings.json")
    from .models import EvidencePlan as _EP, ActivityPlan as _AP, LearningStatePlan as _LSP, PresentationBindingPlan as _PBP
    evidence_plan = _EP(**ep_data) if ep_data else None
    activity_plan = _AP(**ap_data) if ap_data else None
    state_plan = _LSP(**sp_data) if sp_data else None
    activity_bindings = _PBP(**binding_data) if binding_data else None
    deck_plan = build_pptx_deck_plan(
        blueprint, "Chinese", profile.scaffolding_language if profile else "English",
        level or "zero_beginner",
        evidence_plan=evidence_plan, activity_plan=activity_plan, state_plan=state_plan,
        activity_bindings=activity_bindings,
    )
    _wj(project_id, "blueprints/pptx_deck_plan.json", deck_plan.model_dump(mode="json"))
    struct_report = build_pptx_structure_report(deck_plan)
    _wj(project_id, "quality/pptx_structure_report.json", struct_report)
    if struct_report.get("state") == "blocked" and not force:
        raise PermissionError("PPTX structure gate is blocked; shorten or split dense slide content before export")

    prs = _new_master_presentation()
    is_classroom = export_mode == "classroom"
    blank = min(prs.slide_layouts, key=lambda layout: len(layout.placeholders))
    source_slides = {slide.id: slide for slide in blueprint.slides}
    for page_number, deck_slide in enumerate(deck_plan.slides, start=1):
        pptx_slide = prs.slides.add_slide(blank)
        source_slide = source_slides.get(deck_slide.slide_id)
        _render_master_slide(
            pptx_slide, root, source_slide, deck_slide, manifest,
            page_number, len(deck_plan.slides), is_classroom,
        )

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    is_diagnostic = force and is_classroom
    prefix = "HanClassStudio_Diagnostic" if is_diagnostic else "HanClassStudio_Editable"
    export_path = root / "exports" / f"{prefix}_{timestamp}.pptx"
    export_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(export_path)

    quality_state = report.state if report else "warning"
    pptx_report = _build_pptx_quality_report(blueprint, report, force, prs, theme=theme)
    forced_blockers: list[str] = []
    for payload in (alignment_report, readiness_report, binding_report, report.model_dump(mode="json") if report else {}):
        values = payload.get("blocking_reasons", payload.get("blocking", [])) if isinstance(payload, dict) else []
        if isinstance(values, list):
            forced_blockers.extend(str(item) for item in values)
        if force and not values and isinstance(payload, dict) and str(payload.get("state", "")).lower() in {"blocked", "warning"}:
            forced_blockers.append(f"gate state: {payload.get('state')}")
    write_json(project_id, "quality/pptx_quality_report.json", pptx_report)
    write_json(
        project_id,
        "exports/pptx_export_manifest.json",
        {
            "schema": "hanclassstudio.pptx_export_manifest.v1",
            "project_id": project_id,
            "filename": export_path.name,
            "created_at": datetime.now().isoformat(timespec="seconds"),
            "export_type": "pptx_editable",
            "editable": True,
            "forced": force,
            "diagnostic": is_diagnostic,
            "quality_state": quality_state,
            "evidence_alignment_state": alignment_report.get("state") if isinstance(alignment_report, dict) else None,
            "presentation_readiness_state": readiness_report.get("state") if isinstance(readiness_report, dict) else None,
            "forced_blockers": forced_blockers if force else [],
            "force_confirmation": "explicit force=true request" if force else None,
            "interaction_policy": "classroom_static_activity",
            "presentation_theme": {
                "theme_id": theme.theme_id,
                "version": theme.version,
                "source": theme.source,
            },
            "source_artifacts": {
                "spec_lock": bool(spec_lock),
                "interaction_plan": bool(interaction_plan),
                "media_plan": bool(media_plan),
                "activity_bindings": bool(activity_bindings and activity_bindings.bindings),
                "asset_manifest": bool(manifest.model_dump(mode="json")),
            },
        },
    )
    return export_path


def _new_master_presentation() -> Presentation:
    repository_root = Path(__file__).resolve().parents[4]
    master_path = repository_root / PROFILE.source
    prs = Presentation(master_path) if master_path.is_file() else Presentation()
    for slide_id in list(prs.slides._sldIdLst):
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)
    prs.slide_width = Inches(PROFILE.slide_width)
    prs.slide_height = Inches(PROFILE.slide_height)
    return prs


def _render_slide(slide, root: Path, slide_model, manifest: AssetManifest, is_classroom: bool = False) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    slide_type = slide_model.slide_type if slide_model.slide_type in SUPPORTED_SLIDE_TYPES else "GenericContentSlide"
    if not is_classroom:
        _add_label(slide, slide_type)
    if slide_type == "CoverSlide":
        _add_text(slide, slide_model.title, 0.8, 1.55, 7.6, 1.35, 42, bold=True, color=INK)
        _add_text(slide, _content_text(slide_model), 0.85, 3.05, 6.9, 1.2, 21, color=MUTED)
        _add_media_card(slide, root, slide_model, manifest, 8.4, 1.35, 4.1, 3.4, is_classroom)
    elif slide_type == "ObjectiveSlide":
        _add_title(slide, slide_model.title or "学习目标")
        objectives = [block.text for block in slide_model.content_blocks if block.text] or ["完成本课学习任务"]
        for index, text in enumerate(objectives[:4]):
            _add_card(slide, str(index + 1), text, 0.85 + (index % 2) * 6.05, 1.75 + (index // 2) * 1.85, 5.55, 1.45)
    else:
        _add_title(slide, slide_model.title or "课堂活动")
        _add_content_blocks(slide, slide_model)
        _add_media_card(slide, root, slide_model, manifest, 8.55, 1.25, 3.9, 2.35, is_classroom)

    y = 4.0 if slide_model.components else 6.55
    for component in slide_model.components[:3]:
        y = _render_component(slide, component, y)
    if not is_classroom:
        _add_footer(slide, "Editable PPTX export · HTML interactions are converted to classroom static activity pages")


def _render_component(slide, component, y: float) -> float:
    data = component.data or {}
    title = component.title or component.component_type
    if component.component_type == "AudioButton":
        _add_activity_box(slide, f"Audio · {title}", data.get("audio_text") or data.get("label") or "Audio prompt", 0.85, y, 3.7, 1.1, icon="Audio")
        return y + 1.25
    if component.component_type == "VocabularyFlipCard":
        items = _list(data.get("items"))[:4]
        if not items:
            _add_activity_box(slide, title, "Vocabulary cards need items.", 0.85, y, 5.0, 1.0)
            return y + 1.15
        for index, item in enumerate(items):
            x = 0.85 + index * 3.05
            body = "\n".join(
                part
                for part in [
                    item.get("word", ""),
                    item.get("pinyin", ""),
                    item.get("meaning", ""),
                    item.get("example", ""),
                ]
                if part
            )
            _add_activity_box(slide, title if index == 0 else "词卡", body, x, y, 2.75, 1.85)
        return y + 2.05
    if component.component_type == "SentenceDragBuilder":
        words = [str(word) for word in _list(data.get("words"))]
        answer = " ".join(str(word) for word in _list(data.get("answer")))
        _add_activity_box(slide, title, "排序词块: " + "  /  ".join(words), 0.85, y, 7.3, 1.05)
        _add_activity_box(slide, "答案提示", answer or "Teacher answer area", 8.45, y, 3.9, 1.05)
        return y + 1.25
    if component.component_type == "ListenAndChoose":
        choices = "\n".join(f"{i + 1}. {choice}" for i, choice in enumerate(_list(data.get("choices"))))
        body = "\n".join(part for part in [data.get("audio_text", ""), choices, f"Answer: {data.get('answer', '')}"] if part)
        _add_activity_box(slide, title, body or "Listen and choose activity", 0.85, y, 6.4, 1.55)
        return y + 1.75
    if component.component_type == "MatchGame":
        pairs = [pair for pair in _list(data.get("pairs")) if isinstance(pair, dict)]
        body = "\n".join(f"{pair.get('left', '')}   ⟷   {pair.get('right', '')}" for pair in pairs) or "Match pairs"
        _add_activity_box(slide, title, body, 0.85, y, 6.4, 1.55)
        return y + 1.75
    if component.component_type == "CharacterFormation":
        parts = " + ".join(str(part) for part in _list(data.get("parts")))
        character = str(data.get("character", "字"))
        body = f"{parts}  =  {character}\n{data.get('explanation', '')}".strip()
        _add_activity_box(slide, title, body, 0.85, y, 6.4, 1.55)
        return y + 1.75
    if component.component_type == "ClassroomGame":
        _add_activity_box(slide, title, json.dumps(data, ensure_ascii=False, indent=2), 0.85, y, 6.4, 1.55)
        return y + 1.75
    _add_activity_box(slide, f"{title} · static fallback", "Unsupported component rendered as editable activity notes.", 0.85, y, 6.4, 1.1)
    return y + 1.25


def _add_title(slide, title: str) -> None:
    _add_text(slide, title, 0.85, 0.62, 8.8, 0.68, 31, bold=True, color=INK)


def _add_content_blocks(slide, slide_model) -> None:
    text = _content_text(slide_model)
    if text:
        _add_text(slide, text, 0.9, 1.35, 7.0, 2.3, 22, color=INK)


def _add_media_card(slide, root: Path, slide_model, manifest: AssetManifest, x: float, y: float, w: float, h: float, is_classroom: bool = False) -> None:
    path = _image_path(root, slide_model.media_requirements.image_key, manifest)
    if path and path.suffix.lower() in {".png", ".jpg", ".jpeg"}:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        return
    if path and path.suffix.lower() == ".svg":
        png = _rasterize_svg(path)
        if png is not None and png.exists():
            slide.shapes.add_picture(str(png), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
            return
    if is_classroom:
        # Classroom: hide prompt text, use minimal placeholder or nothing
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = MINT
        shape.line.color.rgb = LINE
        shape.line.fill.background()
        return
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MINT
    shape.line.color.rgb = LINE
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = slide_model.media_requirements.image_prompt or "Image placeholder"
    run.font.size = Pt(16)
    run.font.color.rgb = TEAL
    run.font.bold = True


def _add_label(slide, text: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(0.25), Inches(2.35), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MINT
    shape.line.color.rgb = MINT
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = TEAL


def _add_card(slide, label: str, body: str, x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = LINE
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = f"{label}. {body}"
    run.font.size = Pt(19)
    run.font.color.rgb = INK


def _add_activity_box(slide, title: str, body: str, x: float, y: float, w: float, h: float, icon: str = "") -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = LINE
    frame = shape.text_frame
    frame.clear()
    title_p = frame.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = f"{icon} {title}".strip()
    title_run.font.size = Pt(14)
    title_run.font.bold = True
    title_run.font.color.rgb = TEAL
    body_p = frame.add_paragraph()
    body_run = body_p.add_run()
    body_run.text = body or "Static classroom activity"
    body_run.font.size = Pt(15)
    body_run.font.color.rgb = INK


def _add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int, bold: bool = False, color: RGBColor = INK, center: bool = False, font_name: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    if center:
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = PROFILE.latin_font
    if font_name:
        properties = run._r.get_or_add_rPr()
        east_asian = OxmlElement("a:ea")
        east_asian.set("typeface", font_name)
        east_asian.set("panose", "020B0503020204020204")
        east_asian.set("pitchFamily", "34")
        east_asian.set("charset", "-122")
        properties.append(east_asian)


def _add_footer(slide, text: str) -> None:
    _add_text(slide, text, 0.85, 6.93, 11.6, 0.28, 9, color=MUTED)


# ── PPT-master-derived editable renderer ──

def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _render_master_slide(slide, root: Path, source, deck, manifest: AssetManifest, page: int, total: int, is_classroom: bool) -> None:
    """Semantic archetype → master-derived recipe → editable PPTX objects."""
    layout = deck.traditional_layout if deck.traditional_layout in RECIPES else "generic_content"
    _master_background(slide, page, total)
    if source is None:
        _master_generic(slide, None, deck)
    elif layout == "cover_title":
        _master_cover(slide, root, source, deck, manifest)
    elif layout == "objectives_cards":
        _master_objectives(slide, source, deck)
    elif source.slide_type == "PhoneticsSlide":
        _master_phonetics(slide, root, source, manifest)
    elif layout == "two_card_contrast":
        _master_contrast(slide, root, source, deck, manifest)
    elif layout == "listen_choose":
        _master_listen(slide, source, deck)
    elif layout == "dialogue_bubbles":
        _master_dialogue(slide, source, deck)
    elif layout == "match_pairs":
        _master_match(slide, source, deck)
    elif layout == "summary_cards":
        _master_summary(slide, source, deck)
    elif layout == "single_item_focus" and _vocabulary_items(source):
        _master_vocabulary(slide, source, deck)
    elif source.media_requirements.image_key:
        _master_scene(slide, root, source, deck, manifest)
    else:
        _master_generic(slide, source, deck)
    if not is_classroom:
        _add_text(slide, "Editable PPTX · HanClassStudio", 0.72, 7.05, 3.0, 0.2, 10, color=_rgb(PROFILE.muted))
    if deck.speaker_notes:
        notes = slide.notes_slide.notes_text_frame
        notes.clear()
        notes.paragraphs[0].add_run().text = "\n".join(deck.speaker_notes)


def _master_background(slide, page: int, total: int) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(PROFILE.background)
    for x, y, size, color, transparency in [
        (0.0, 0.0, 1.05, PROFILE.primary, 70),
        (0.28, 0.0, 0.5, PROFILE.accent, 58),
        (12.3, 6.48, 1.02, PROFILE.primary, 62),
        (12.02, 6.91, 0.5, PROFILE.accent, 45),
    ]:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        shape.fill.solid(); shape.fill.fore_color.rgb = _rgb(color); shape.fill.transparency = transparency
        shape.line.fill.background()
    _add_text(slide, f"{page:02d} / {total:02d}", 11.55, 0.28, 1.0, 0.28, 10, bold=True, color=_rgb(PROFILE.primary), center=True)


def _master_header(slide, title: str, instruction: str = "") -> None:
    _add_text(slide, title, 0.72, 0.42, 8.7, 0.55, PROFILE.title_size, bold=True, color=_rgb(PROFILE.ink), font_name=PROFILE.heading_font)
    if instruction:
        _add_text(slide, instruction, 0.74, 0.94, 10.8, 0.38, 13, color=_rgb(PROFILE.muted), font_name=PROFILE.latin_font)


def _master_cover(slide, root: Path, source, deck, manifest: AssetManifest) -> None:
    _add_text(slide, source.title, 0.78, 1.62, 6.0, 1.15, 48, bold=True, color=_rgb(PROFILE.ink), font_name=PROFILE.chinese_font)
    block = source.content_blocks[0] if source.content_blocks else None
    if block:
        pinyin, meaning = _split_scaffold(block.scaffolding_text)
        _add_text(slide, pinyin, 0.82, 3.03, 5.7, 0.48, PROFILE.pinyin_size, color=_rgb(PROFILE.primary), font_name=PROFILE.latin_font)
        _add_text(slide, meaning, 0.82, 3.58, 5.7, 0.5, PROFILE.meaning_size, color=_rgb(PROFILE.muted), font_name=PROFILE.latin_font)
    _add_text(slide, source.title, 0.82, 5.9, 5.8, 0.4, 16, bold=True, color=_rgb(PROFILE.accent), font_name=PROFILE.heading_font)
    _master_picture(slide, _image_path(root, source.media_requirements.image_key, manifest), 7.15, 0.75, 5.45, 5.95)


def _master_objectives(slide, source, deck) -> None:
    _master_header(slide, source.title or "学习目标", "By the end of this lesson")
    objectives = [block.text for block in source.content_blocks if block.text] or deck.target_text.splitlines()
    for index, text in enumerate(objectives[:4]):
        y = 1.6 + index * 1.22
        _master_card(slide, 1.15, y, 11.0, 0.88, PROFILE.background_alt)
        _add_text(slide, str(index + 1), 1.4, y + 0.17, 0.45, 0.4, 22, bold=True, color=_rgb(PROFILE.accent), center=True)
        _add_text(slide, text, 2.05, y + 0.13, 9.55, 0.5, 20, color=_rgb(PROFILE.ink), font_name=PROFILE.heading_font)


def _master_vocabulary(slide, source, deck) -> None:
    _master_header(slide, source.title, "Look, read, and say")
    items = _vocabulary_items(source)[:6]
    if len(items) == 2:
        for index, item in enumerate(items):
            x, y = 0.85 + index * 6.08, 1.55
            _master_card(slide, x, y, 5.58, 4.35, "FFFFFF")
            _add_text(slide, item.get("word", ""), x + 0.3, y + 0.48, 4.98, 0.82, 42, bold=True, color=_rgb(PROFILE.ink), center=True, font_name=PROFILE.chinese_font)
            _add_text(slide, item.get("pinyin", ""), x + 0.3, y + 1.55, 4.98, 0.42, 22, color=_rgb(PROFILE.primary), center=True)
            _add_text(slide, item.get("meaning", ""), x + 0.35, y + 2.24, 4.88, 0.5, 18, color=_rgb(PROFILE.ink), center=True)
            usage_context = str(item.get("usage_context", "")).strip()
            if usage_context:
                _add_text(slide, usage_context, x + 0.52, y + 3.15, 4.54, 0.65, 14, color=_rgb(PROFILE.muted), center=True)
        return
    cols = 3
    for index, item in enumerate(items):
        row, col = divmod(index, cols)
        x, y = 0.85 + col * 4.08, 1.55 + row * 2.38
        _master_card(slide, x, y, 3.65, 1.95, "FFFFFF")
        _add_text(slide, item.get("word", ""), x + 0.2, y + 0.24, 3.25, 0.62, 34, bold=True, color=_rgb(PROFILE.ink), center=True, font_name=PROFILE.chinese_font)
        _add_text(slide, item.get("pinyin", ""), x + 0.2, y + 0.95, 3.25, 0.34, 18, color=_rgb(PROFILE.primary), center=True)
        _add_text(slide, item.get("meaning", ""), x + 0.2, y + 1.38, 3.25, 0.3, 14, color=_rgb(PROFILE.muted), center=True)


def _master_scene(slide, root: Path, source, deck, manifest: AssetManifest) -> None:
    _master_header(slide, source.title, "Look and say")
    _master_picture(slide, _image_path(root, source.media_requirements.image_key, manifest), 6.45, 1.42, 5.95, 4.95)
    if source.content_blocks:
        block = source.content_blocks[0]
        pinyin, meaning = _split_scaffold(block.scaffolding_text)
        _add_text(slide, block.text, 0.85, 2.0, 5.05, 0.9, 44, bold=True, color=_rgb(PROFILE.ink), font_name=PROFILE.chinese_font)
        _add_text(slide, pinyin, 0.87, 3.12, 4.9, 0.45, PROFILE.pinyin_size, color=_rgb(PROFILE.primary))
        _add_text(slide, meaning, 0.87, 3.75, 4.9, 0.55, PROFILE.meaning_size, color=_rgb(PROFILE.muted))


def _master_phonetics(slide, root: Path, source, manifest: AssetManifest) -> None:
    """Render source-first phonetics without text/image collisions."""
    _master_header(slide, source.title, "Read the examples; use the visual to explain the sound pattern.")
    if "声母" in source.title and "韵母" in source.title:
        block = source.content_blocks[0] if source.content_blocks else None
        _phonetics_copy(slide, block)
        _draw_syllable_structure(slide)
        return
    if "声调位置" in source.title:
        block = source.content_blocks[0] if source.content_blocks else None
        _phonetics_copy(slide, block)
        _draw_tone_mark_examples(slide)
        return
    if "声调" in source.title and "变调" not in source.title:
        block = source.content_blocks[0] if source.content_blocks else None
        _phonetics_copy(slide, block)
        _draw_tone_contours(slide)
        return

    image = _image_path(root, source.media_requirements.image_key, manifest)
    if image:
        block = source.content_blocks[0] if source.content_blocks else None
        target = block.text if block else ""
        scaffold = block.scaffolding_text if block else ""
        line_count = max(1, len(target.splitlines()))
        target_size = 38 if line_count == 1 else 31
        _add_text(
            slide, target, 0.88, 1.72, 5.15, 2.05, target_size,
            bold=True, color=_rgb(PROFILE.ink), font_name=PROFILE.chinese_font,
        )
        _add_text(
            slide, scaffold, 0.9, 4.08, 5.05, 1.25, 17,
            color=_rgb(PROFILE.primary), font_name=PROFILE.latin_font,
        )
        _master_picture(slide, image, 6.55, 1.45, 5.7, 4.85)
        return

    for index, block in enumerate(source.content_blocks[:4]):
        row, col = divmod(index, 2)
        x, y = 0.85 + col * 6.08, 1.52 + row * 2.45
        _master_card(slide, x, y, 5.58, 2.03, "FFFFFF" if index % 2 == 0 else PROFILE.background_alt)
        _add_text(
            slide, block.text, x + 0.25, y + 0.25, 5.08, 0.62, 23,
            bold=True, color=_rgb(PROFILE.ink), font_name=PROFILE.chinese_font,
        )
        _add_text(
            slide, block.scaffolding_text, x + 0.25, y + 0.98, 5.08, 0.78, 14,
            color=_rgb(PROFILE.muted), font_name=PROFILE.latin_font,
        )


def _phonetics_copy(slide, block) -> None:
    target = block.text if block else ""
    scaffold = block.scaffolding_text if block else ""
    line_count = max(1, len(target.splitlines()))
    _add_text(
        slide, target, 0.88, 1.72, 5.15, 2.05, 38 if line_count == 1 else 31,
        bold=True, color=_rgb(PROFILE.ink), font_name=PROFILE.chinese_font,
    )
    _add_text(
        slide, scaffold, 0.9, 4.08, 5.05, 1.25, 17,
        color=_rgb(PROFILE.primary), font_name=PROFILE.latin_font,
    )


def _draw_syllable_structure(slide) -> None:
    _master_card(slide, 6.55, 1.55, 5.7, 4.72, "FFFFFF")
    _add_text(slide, "声母", 7.25, 2.45, 1.55, 0.58, 29, bold=True, color=_rgb(PROFILE.ink), center=True)
    _add_text(slide, "+", 8.9, 2.5, 0.55, 0.48, 28, bold=True, color=_rgb(PROFILE.accent), center=True)
    _add_text(slide, "韵母", 9.55, 2.45, 1.55, 0.58, 29, bold=True, color=_rgb(PROFILE.ink), center=True)
    _add_text(slide, "initial", 7.25, 3.15, 1.55, 0.32, 14, color=_rgb(PROFILE.primary), center=True)
    _add_text(slide, "final", 9.55, 3.15, 1.55, 0.32, 14, color=_rgb(PROFILE.primary), center=True)
    _add_text(slide, "↓", 8.98, 3.72, 0.5, 0.45, 25, bold=True, color=_rgb(PROFILE.accent), center=True)
    _add_text(slide, "一个音节", 7.7, 4.35, 3.0, 0.62, 27, bold=True, color=_rgb(PROFILE.primary), center=True)
    _add_text(slide, "one syllable", 7.7, 5.03, 3.0, 0.35, 14, color=_rgb(PROFILE.muted), center=True)


def _draw_tone_contours(slide) -> None:
    _master_card(slide, 6.55, 1.55, 5.7, 4.72, "FFFFFF")
    starts = [7.05, 8.35, 9.65, 10.95]
    segments = [
        [(0.0, 0.45, 0.95, 0.45)],
        [(0.0, 1.95, 0.95, 0.35)],
        [(0.0, 0.85, 0.48, 1.85), (0.48, 1.85, 0.95, 0.7)],
        [(0.0, 0.35, 0.95, 2.05)],
    ]
    for index, x in enumerate(starts):
        for x1, y1, x2, y2 in segments[index]:
            line = slide.shapes.add_connector(1, Inches(x + x1), Inches(2.2 + y1), Inches(x + x2), Inches(2.2 + y2))
            line.line.color.rgb = _rgb(PROFILE.primary)
            line.line.width = Pt(4)
    for index, (x, label) in enumerate(zip(starts, ("1st", "2nd", "3rd", "4th"))):
        _add_text(slide, label, x, 4.72, 0.95, 0.34, 15, bold=True, color=_rgb(PROFILE.ink), center=True)
    _add_text(slide, "level       rising       dipping       falling", 6.95, 5.25, 4.95, 0.35, 14, color=_rgb(PROFILE.muted), center=True)


def _draw_tone_mark_examples(slide) -> None:
    _master_card(slide, 6.55, 1.55, 5.7, 4.72, "FFFFFF")
    examples = [
        ("mā", "one vowel → mark it"),
        ("hǎo · zuò", "mark the main vowel"),
        ("liú · guǐ", "iu / ui → mark the second"),
    ]
    for index, (example, rule) in enumerate(examples):
        y = 1.92 + index * 1.34
        _master_card(slide, 6.92, y, 4.95, 1.02, PROFILE.background_alt if index % 2 else "FFFFFF")
        _add_text(slide, example, 7.15, y + 0.14, 1.75, 0.48, 25, bold=True, color=_rgb(PROFILE.ink), center=True)
        _add_text(slide, rule, 9.02, y + 0.2, 2.55, 0.4, 14, color=_rgb(PROFILE.muted), center=True)


def _master_contrast(slide, root: Path, source, deck, manifest: AssetManifest) -> None:
    _master_header(slide, source.title or "你好 / 您好", "Choose the greeting for the relationship")
    image = _image_path(root, source.media_requirements.image_key, manifest)
    if image:
        _master_picture(slide, image, 8.05, 1.48, 4.55, 4.65)
        card_width, gap = 3.25, 0.34
    else:
        card_width, gap = 5.65, 0.42
    blocks = source.content_blocks[:2]
    defaults = [("你好！", "nǐ hǎo", "friends and classmates"), ("您好！", "nín hǎo", "teachers and elders")]
    for index in range(2):
        x = 0.78 + index * (card_width + gap)
        block = blocks[index] if index < len(blocks) else None
        hanzi = block.text if block else defaults[index][0]
        pinyin, meaning = _split_scaffold(block.scaffolding_text) if block else defaults[index][1:]
        color = PROFILE.warm if index == 0 else PROFILE.background_alt
        _master_card(slide, x, 1.58, card_width, 4.42, color)
        _add_text(slide, "同伴 · PEER" if index == 0 else "老师 · TEACHER", x + 0.24, 1.88, card_width - 0.48, 0.3, 13, bold=True, color=_rgb(PROFILE.primary), center=True)
        _add_text(slide, hanzi, x + 0.2, 2.5, card_width - 0.4, 0.8, 38, bold=True, color=_rgb(PROFILE.ink), center=True, font_name=PROFILE.chinese_font)
        _add_text(slide, pinyin, x + 0.2, 3.48, card_width - 0.4, 0.38, 20, color=_rgb(PROFILE.primary), center=True)
        _add_text(slide, meaning, x + 0.3, 4.2, card_width - 0.6, 0.72, 16, color=_rgb(PROFILE.muted), center=True)


def _master_dialogue(slide, source, deck) -> None:
    """Keep each textbook line, pinyin, and scaffold meaning in one readable turn card."""
    _master_header(slide, source.title, "Read the source dialogue, then perform it with a partner.")
    blocks = source.content_blocks[:4]
    for index, block in enumerate(blocks):
        row, col = divmod(index, 2)
        x, y = 0.85 + col * 6.08, 1.55 + row * 2.4
        _master_card(slide, x, y, 5.58, 1.92, "FFFFFF" if index % 2 == 0 else PROFILE.background_alt)
        pinyin, meaning = _split_scaffold(block.scaffolding_text)
        _add_text(
            slide, block.text, x + 0.28, y + 0.22, 5.02, 0.52, 25,
            bold=True, color=_rgb(PROFILE.ink), font_name=PROFILE.chinese_font,
        )
        _add_text(slide, pinyin, x + 0.28, y + 0.86, 5.02, 0.32, 16, color=_rgb(PROFILE.primary))
        _add_text(slide, meaning, x + 0.28, y + 1.3, 5.02, 0.3, 14, color=_rgb(PROFILE.muted))


def _master_listen(slide, source, deck) -> None:
    instruction = _first_scaffold(source) or "Listen and choose the greeting you hear."
    _master_header(slide, source.title or "听一听，选一选", instruction)
    component = next((c for c in source.components if c.component_type == "ListenAndChoose"), None)
    data = component.data if component else {}
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.92), Inches(1.62), Inches(0.72), Inches(0.72))
    badge.fill.solid(); badge.fill.fore_color.rgb = _rgb(PROFILE.primary); badge.line.fill.background()
    _add_text(slide, "▶", 1.08, 1.78, 0.38, 0.28, 18, bold=True, color=RGBColor(255, 255, 255), center=True)
    _add_text(slide, "播放本地普通话音频", 1.82, 1.72, 4.8, 0.45, 20, bold=True, color=_rgb(PROFILE.ink), font_name=PROFILE.heading_font)
    choices = [str(choice) for choice in data.get("choices", [])][:4]
    card_width = (11.45 - max(0, len(choices) - 1) * 0.35) / max(1, len(choices))
    for index, choice in enumerate(choices):
        x = 0.95 + index * (card_width + 0.35)
        _master_card(slide, x, 3.05, card_width, 1.78, "FFFFFF")
        _add_text(slide, chr(65 + index), x + 0.18, 3.26, 0.45, 0.42, 18, bold=True, color=_rgb(PROFILE.accent), center=True)
        _add_text(slide, choice, x + 0.62, 3.48, card_width - 0.82, 0.62, 32, bold=True, color=_rgb(PROFILE.ink), center=True, font_name=PROFILE.chinese_font)


def _master_match(slide, source, deck) -> None:
    instruction = _first_scaffold(source)
    if not instruction and source.content_blocks:
        instruction = source.content_blocks[0].text
    _master_header(slide, source.title or "连一连", instruction or "Match each textbook line with its response.")
    component = next((c for c in source.components if c.component_type == "MatchGame"), None)
    pairs = component.data.get("pairs", []) if component else []
    right_values = [str(pair.get("right", "")) for pair in pairs[:5]]
    if len(right_values) > 1:
        right_values = right_values[1:] + right_values[:1]
    _add_text(slide, "Textbook line", 1.0, 1.35, 4.65, 0.3, 13, bold=True, color=_rgb(PROFILE.muted), center=True)
    _add_text(slide, "Response", 7.65, 1.35, 4.65, 0.3, 13, bold=True, color=_rgb(PROFILE.muted), center=True)
    for index, pair in enumerate(pairs[:5]):
        y = 1.72 + index * 1.0
        _master_card(slide, 1.0, y, 4.65, 0.7, "FFFFFF")
        _master_card(slide, 7.65, y, 4.65, 0.7, PROFILE.background_alt)
        _add_text(slide, str(pair.get("left", "")), 1.25, y + 0.12, 4.15, 0.38, 23, bold=True, color=_rgb(PROFILE.ink), center=True, font_name=PROFILE.chinese_font)
        _add_text(slide, right_values[index], 7.9, y + 0.13, 4.15, 0.36, 20, color=_rgb(PROFILE.ink), center=True, font_name=PROFILE.chinese_font)
        _add_text(slide, str(index + 1), 5.95, y + 0.15, 0.35, 0.3, 13, bold=True, color=_rgb(PROFILE.primary), center=True)
        _add_text(slide, chr(65 + index), 7.0, y + 0.15, 0.35, 0.3, 13, bold=True, color=_rgb(PROFILE.accent), center=True)


def _master_summary(slide, source, deck) -> None:
    _master_header(slide, source.title or "今天会说", "Review and choose the greeting that fits")
    text = " · ".join(block.text for block in source.content_blocks if block.text) or deck.target_text
    items = [item.strip() for item in text.split("·") if item.strip()]
    for index, item in enumerate(items[:5]):
        x = 0.85 + (index % 3) * 4.08
        y = 1.72 + (index // 3) * 2.0
        _master_card(slide, x, y, 3.65, 1.52, PROFILE.background_alt if index % 2 else "FFFFFF")
        _add_text(slide, item, x + 0.2, y + 0.42, 3.25, 0.56, 28, bold=True, color=_rgb(PROFILE.ink), center=True, font_name=PROFILE.chinese_font)


def _master_generic(slide, source, deck) -> None:
    title = source.title if source else deck.main_focus
    _master_header(slide, title or "课堂活动")
    text = _content_text(source) if source else deck.target_text
    _add_text(slide, text, 1.0, 1.72, 11.2, 4.8, 24, color=_rgb(PROFILE.ink), font_name=PROFILE.heading_font)


def _master_card(slide, x: float, y: float, w: float, h: float, color: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = _rgb(color)
    shape.line.color.rgb = _rgb(PROFILE.line); shape.line.width = Pt(1.2)


def _master_picture(slide, path: Path | None, x: float, y: float, w: float, h: float) -> None:
    _master_card(slide, x - 0.06, y - 0.06, w + 0.12, h + 0.12, "FFFFFF")
    if path and path.suffix.lower() == ".svg":
        svg_path = path
        path = _rasterize_svg(path)
        if path is None:
            raise RuntimeError(f"PPTX export could not rasterize SVG asset: {svg_path.name}")
    if not path or path.suffix.lower() not in {".png", ".jpg", ".jpeg"}:
        return
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    box_ratio = w / h
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if image_ratio > box_ratio:
        crop = (1 - box_ratio / image_ratio) / 2
        picture.crop_left = crop; picture.crop_right = crop
    elif image_ratio < box_ratio:
        crop = (1 - image_ratio / box_ratio) / 2
        picture.crop_top = crop; picture.crop_bottom = crop


def _vocabulary_items(source) -> list[dict]:
    return [item for component in source.components if component.component_type == "VocabularyFlipCard" for item in component.data.get("items", []) if isinstance(item, dict)]


def _split_scaffold(text: str) -> tuple[str, str]:
    parts = [part.strip() for part in text.split("·", 1)]
    return (parts[0], parts[1]) if len(parts) == 2 else (text.strip(), "")


def _first_scaffold(source) -> str:
    return next((block.scaffolding_text for block in source.content_blocks if block.scaffolding_text), "")


# ── New: Traditional Deck Slide Renderer ──

SUPPORTED_DECK_LAYOUTS: set[str] = {
    "cover_title", "objectives_cards", "single_item_focus",
    "two_card_contrast", "listen_choose", "dialogue_bubbles",
    "match_pairs", "summary_cards", "generic_content",
}


def _render_deck_slide(slide, deck_slide, is_classroom: bool = False) -> None:
    """Render a traditional deck slide from a PptxDeckSlide plan."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    layout = deck_slide.traditional_layout if deck_slide.traditional_layout in SUPPORTED_DECK_LAYOUTS else "generic_content"

    if layout == "single_item_focus":
        _render_deck_single_item(slide, deck_slide)
    elif layout == "two_card_contrast":
        _render_deck_two_card_contrast(slide, deck_slide)
    elif layout == "cover_title":
        _render_deck_cover(slide, deck_slide)
    elif layout == "objectives_cards":
        _render_deck_objectives(slide, deck_slide)
    elif layout == "dialogue_bubbles":
        _render_deck_dialogue(slide, deck_slide)
    elif layout == "match_pairs":
        _render_deck_match(slide, deck_slide)
    elif layout == "summary_cards":
        _render_deck_summary(slide, deck_slide)
    else:
        _render_deck_generic(slide, deck_slide)

    if not is_classroom:
        _add_footer(slide, "Editable PPTX · HanClassStudio")
    if deck_slide.speaker_notes:
        notes_slide = slide.notes_slide
        if notes_slide:
            tf = notes_slide.notes_text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.add_run().text = "\n".join(deck_slide.speaker_notes)


def _render_deck_single_item(slide, ds) -> None:
    """Hero target word centered, pronunciation below, scaffold below."""
    # Decorative background
    bg = slide.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(12.3), Inches(6.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    bg.line.fill.background()
    # Hero text
    _add_text(slide, ds.target_text or ds.main_focus, 2.5, 1.2, 8.3, 2.0, 52, bold=True, color=RGBColor(0x33, 0x33, 0x33), center=True)
    # Pronunciation
    if ds.pronunciation:
        _add_text(slide, ds.pronunciation, 2.5, 3.3, 8.3, 0.7, 26, color=RGBColor(0x99, 0x99, 0x99), center=True)
    # Scaffold meaning
    if ds.scaffold_text:
        _add_text(slide, ds.scaffold_text, 2.5, 4.1, 8.3, 0.8, 22, color=RGBColor(0x66, 0x66, 0x66), center=True)
    # Usage context
    if ds.usage_context:
        _add_text(slide, ds.usage_context, 2.5, 5.0, 8.3, 0.6, 18, color=RGBColor(0x99, 0x99, 0x99), center=True)
    # Teacher notes in classroom mode only via speaker notes


def _render_deck_two_card_contrast(slide, ds) -> None:
    """Two scene cards side by side."""
    left_label = ds.target_text.split("/")[0].strip() if "/" in ds.target_text else ""
    right_label = ds.target_text.split("/")[1].strip() if "/" in ds.target_text else ""
    # Left card
    card1 = slide.shapes.add_shape(1, Inches(0.8), Inches(1.0), Inches(5.5), Inches(5.0))
    card1.fill.solid(); card1.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    card1.line.fill.background()
    _add_text(slide, left_label or "你好！", 1.0, 2.0, 5.1, 1.5, 40, bold=True, color=RGBColor(0x33, 0x33, 0x33), center=True)
    _add_text(slide, ds.scaffold_text, 1.0, 3.8, 5.1, 1.0, 20, color=RGBColor(0x66, 0x66, 0x66), center=True)
    # Right card
    card2 = slide.shapes.add_shape(1, Inches(6.8), Inches(1.0), Inches(5.5), Inches(5.0))
    card2.fill.solid(); card2.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    card2.line.fill.background()
    _add_text(slide, right_label or "您好！", 7.0, 2.0, 5.1, 1.5, 40, bold=True, color=RGBColor(0x33, 0x33, 0x33), center=True)
    _add_text(slide, ds.usage_context or ds.scaffold_text, 7.0, 3.8, 5.1, 1.0, 20, color=RGBColor(0x66, 0x66, 0x66), center=True)


def _render_deck_cover(slide, ds) -> None:
    _add_text(slide, ds.target_text or ds.main_focus, 2.0, 2.5, 9.3, 2.0, 48, bold=True, color=RGBColor(0x33, 0x33, 0x33), center=True)


def _render_deck_objectives(slide, ds) -> None:
    _add_text(slide, "学习目标", 1.0, 0.5, 11.3, 0.8, 28, bold=True, color=RGBColor(0x33, 0x33, 0x33))
    lines = ds.target_text.split("\\n")
    y = 1.5
    for line in lines[:4]:
        _add_text(slide, f"• {line}", 1.5, y, 10.3, 0.7, 20, color=RGBColor(0x33, 0x33, 0x33))
        y += 1.0


def _render_deck_dialogue(slide, ds) -> None:
    lines = ds.target_text.split("\\n")
    y = 1.0
    for i, line in enumerate(lines[:6]):
        side = 1.0 if i % 2 == 0 else 4.5
        _add_text(slide, line, side, y, 8.0, 0.7, 20, color=RGBColor(0x33, 0x33, 0x33))
        y += 0.9


def _render_deck_match(slide, ds) -> None:
    lines = ds.target_text.split("\\n")
    y = 1.5
    for line in lines[:8]:
        _add_text(slide, line, 2.0, y, 9.3, 0.6, 18, color=RGBColor(0x33, 0x33, 0x33))
        y += 0.8


def _render_deck_summary(slide, ds) -> None:
    _add_text(slide, "课堂小结", 1.0, 0.5, 11.3, 0.8, 28, bold=True, color=RGBColor(0x33, 0x33, 0x33))
    lines = ds.target_text.split("\\n")
    y = 1.5
    for line in lines[:4]:
        _add_text(slide, f"☐ {line}", 1.5, y, 10.3, 0.7, 20, color=RGBColor(0x33, 0x33, 0x33))
        y += 1.0


def _render_deck_generic(slide, ds) -> None:
    _add_text(slide, ds.main_focus or ds.target_text or "课堂活动", 1.0, 0.5, 11.3, 0.8, 28, bold=True, color=RGBColor(0x33, 0x33, 0x33))
    if ds.target_text:
        _add_text(slide, ds.target_text, 1.0, 1.5, 11.3, 4.0, 20, color=RGBColor(0x33, 0x33, 0x33))


def _content_text(slide_model) -> str:
    lines = []
    for block in slide_model.content_blocks:
        if block.text:
            lines.append(block.text)
        if block.scaffolding_text and not PROVIDER_REQUIRED_PPTX.search(block.scaffolding_text):
            lines.append(block.scaffolding_text)
    return "\n".join(lines)


def _rasterize_svg(svg_path: Path) -> Path | None:
    """Rasterize SVG through the existing PyMuPDF dependency for PPTX embedding."""
    try:
        import fitz

        document = fitz.open(stream=svg_path.read_bytes(), filetype="svg")
        page = document[0]
        rect = page.rect
        if rect.width <= 0 or rect.height <= 0:
            return None
        pixmap = page.get_pixmap(
            matrix=fitz.Matrix(1200 / rect.width, 675 / rect.height),
            alpha=False,
        )
        out = svg_path.with_suffix(".png")
        pixmap.save(out)
        document.close()
        return out
    except Exception:
        return None


def _image_path(root: Path, image_key: str | None, manifest: AssetManifest) -> Path | None:
    if not image_key:
        return None
    for asset in manifest.images:
        if asset.id == image_key:
            path = root / asset.path
            return path if path.exists() else None
    return None


def _build_pptx_quality_report(blueprint: LessonBlueprint, report: QualityReport | None, force: bool, presentation=None, theme=None) -> dict[str, Any]:
    warnings = [
        "HTML interactions were converted to editable classroom static activity pages.",
        "Audio is represented as text prompts or labels; real audio is not embedded.",
    ]
    if force:
        warnings.append("PPTX export was forced.")
    if not report:
        warnings.append("Source quality_report.json was missing.")
    structural = _pptx_structural_findings(presentation) if presentation else {"off_slide": [], "small_text": []}
    return {
        "schema": "hanclassstudio.pptx_quality_report.v1",
        "state": "warning" if warnings else "pass",
        "blocking": [],
        "warnings": warnings,
        "master_profile": {
            "source": PROFILE.source,
            "slide_size_inches": [PROFILE.slide_width, PROFILE.slide_height],
            "heading_font": PROFILE.heading_font,
            "chinese_font": PROFILE.chinese_font,
            "latin_font": PROFILE.latin_font,
        },
        "presentation_theme": {
            "theme_id": theme.theme_id,
            "version": theme.version,
        } if theme else None,
        "off_slide_objects": structural["off_slide"],
        "text_below_minimum": structural["small_text"],
        "passed": [
            "pptx_file_created",
            f"slide_count:{len(blueprint.slides)}",
            "editable_shapes_created",
        ],
        "source_quality_state": report.state if report else None,
    }


def _pptx_structural_findings(presentation) -> dict[str, list[str]]:
    off_slide: list[str] = []
    small_text: list[str] = []
    width, height = presentation.slide_width, presentation.slide_height
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > width or shape.top + shape.height > height:
                off_slide.append(f"S{slide_index}:shape{shape_index}")
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip() and run.font.size and run.font.size.pt < 10:
                            small_text.append(f"S{slide_index}:{run.text[:24]}:{run.font.size.pt:g}pt")
    return {"off_slide": off_slide, "small_text": small_text}


def _list(value) -> list:
    return value if isinstance(value, list) else []
