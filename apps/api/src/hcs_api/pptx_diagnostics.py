"""Development-only PPTX slide screenshots and contact sheets."""

from __future__ import annotations

import argparse
import shutil
import subprocess
import tempfile
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation


def render_pptx_slides(pptx_path: Path, output_dir: Path) -> list[Path]:
    """Render each slide through macOS Quick Look without flattening the source deck."""
    qlmanage = shutil.which("qlmanage")
    if not qlmanage:
        raise RuntimeError("qlmanage is required for PPTX diagnostic screenshots")
    output_dir.mkdir(parents=True, exist_ok=True)
    slide_count = len(Presentation(pptx_path).slides)
    rendered: list[Path] = []
    with tempfile.TemporaryDirectory(prefix="hcs-pptx-review-") as temp:
        temp_root = Path(temp)
        for index in range(slide_count):
            presentation = Presentation(pptx_path)
            for remove_index in reversed(range(slide_count)):
                if remove_index == index:
                    continue
                slide_id = presentation.slides._sldIdLst[remove_index]
                presentation.part.drop_rel(slide_id.rId)
                presentation.slides._sldIdLst.remove(slide_id)
            single = temp_root / f"slide-{index + 1:02d}.pptx"
            presentation.save(single)
            subprocess.run(
                [qlmanage, "-t", "-s", "1600", "-o", str(temp_root), str(single)],
                check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
            )
            source = temp_root / f"{single.name}.png"
            target = output_dir / f"slide-{index + 1:02d}.png"
            shutil.move(source, target)
            rendered.append(target)
    return rendered


def build_contact_sheet(images: list[Path], output: Path, columns: int = 3) -> Path:
    cell_width, cell_height = 426, 264
    rows = (len(images) + columns - 1) // columns
    sheet = Image.new("RGB", (columns * cell_width, rows * cell_height), "#DDE7EF")
    for index, path in enumerate(images):
        image = Image.open(path).convert("RGB")
        image.thumbnail((400, 225))
        cell = Image.new("RGB", (cell_width, cell_height), "white")
        cell.paste(image, ((cell_width - image.width) // 2, 12))
        ImageDraw.Draw(cell).text((12, 242), f"Slide {index + 1}", fill="#26374A")
        sheet.paste(cell, ((index % columns) * cell_width, (index // columns) * cell_height))
    output.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output)
    return output


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("output_dir", type=Path)
    args = parser.parse_args()
    paths = render_pptx_slides(args.pptx, args.output_dir)
    print(build_contact_sheet(paths, args.output_dir / "contact-sheet.png"))
