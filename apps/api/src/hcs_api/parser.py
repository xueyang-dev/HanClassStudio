from __future__ import annotations

import unicodedata
from pathlib import Path

from .models import ImageBlock, SourceMaterial, SourcePage, TextBlock
from .source_understanding import OCRPolicy, PageInput, run_source_understanding

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg"}
TITLE_BLOCK_TYPES = {"title", "unit_title", "lesson_title", "heading"}


def parse_source(
    file_path: Path,
    project_root: Path,
    original_filename: str,
    force_engine: str | None = None,
) -> SourceMaterial:
    """Parse an uploaded teaching file and run the OCR / Source Document
    Understanding layer. Returns a SourceMaterial whose ``source_analysis`` holds
    the structured, evidence-preserving Source Evidence Model.

    Supported inputs: PPTX, PDF, and scanned page images (PNG/JPEG).

    ``force_engine`` (optional) overrides the OCR engine selector for scanned /
    PDF pages so a teacher can re-run recognition with a specific engine (e.g.
    ``"paddle_ocr"`` or ``"tesseract"``). It is ignored for PPTX, which always
    uses its reliable native text layer.
    """
    suffix = file_path.suffix.lower()
    if suffix == ".pptx":
        return parse_pptx(file_path, project_root, original_filename, force_engine)
    if suffix == ".pdf":
        return parse_pdf(file_path, project_root, original_filename, force_engine)
    if suffix in IMAGE_SUFFIXES:
        return parse_image(file_path, project_root, original_filename, force_engine)
    raise ValueError("Only PPTX, PDF and PNG/JPEG files are supported in v0.1+")


# ── PPTX ─────────────────────────────────────────────────────────────────────

def parse_pptx(
    file_path: Path,
    project_root: Path,
    original_filename: str,
    force_engine: str | None = None,
) -> SourceMaterial:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(str(file_path))
    slide_w = _emu_to_points(getattr(presentation, "slide_width", None)) or 1.0
    slide_h = _emu_to_points(getattr(presentation, "slide_height", None)) or 1.0

    page_inputs: list[PageInput] = []
    legacy_images: list[list[tuple[str, list[float] | None]]] = []
    titles: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        title = ""
        native_blocks: list[tuple[str, list[float] | None, str]] = []
        embedded: list[tuple[str, list[float] | None]] = []
        image_dir = project_root / "assets" / "images"
        image_dir.mkdir(parents=True, exist_ok=True)

        if slide.shapes.title and getattr(slide.shapes.title, "text", "").strip():
            title = slide.shapes.title.text.strip()

        image_count = 0
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(
                    p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                ).strip()
                if text:
                    kind = "title" if (not title or text == title) else "body"
                    if not title:
                        title = text.splitlines()[0]
                    bbox = _shape_bbox_norm(shape, slide_w, slide_h)
                    native_blocks.append((text, bbox, kind))

            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, "image"):
                image_count += 1
                image = shape.image
                ext = image.ext or "png"
                filename = f"source_p{slide_index}_{image_count}.{ext}"
                (image_dir / filename).write_bytes(image.blob)
                rel = f"assets/images/{filename}"
                embedded.append((rel, _shape_bbox_norm(shape, slide_w, slide_h)))

        whole_slide_image = next(
            (rel for rel, bbox in embedded if bbox and _bbox_area(bbox) >= 0.8),
            None,
        ) if not native_blocks else None
        image_path = project_root / whole_slide_image if whole_slide_image else None
        image_size = _image_size(image_path) if image_path else None
        page_inputs.append(PageInput(
            page_number=slide_index,
            width=image_size[0] if image_size else slide_w,
            height=image_size[1] if image_size else slide_h,
            has_native_text=bool(native_blocks),
            native_blocks=native_blocks,
            image=str(image_path) if image_path else None,
            embedded_images=embedded,
            is_whole_image=bool(image_path),
        ))
        legacy_images.append(embedded)
        titles.append(title or f"第 {slide_index} 页")

    return _assemble(page_inputs, legacy_images, titles, "pptx", original_filename, project_root, force_engine)


# ── PDF ──────────────────────────────────────────────────────────────────────

def parse_pdf(
    file_path: Path,
    project_root: Path,
    original_filename: str,
    force_engine: str | None = None,
) -> SourceMaterial:
    try:
        import fitz
    except ImportError as exc:
        raise ValueError("PyMuPDF is required for PDF parsing") from exc

    document = fitz.open(file_path)
    render_dir = project_root / "assets" / "ocr_render"
    crop_dir = project_root / "assets" / "ocr_crops"
    render_dir.mkdir(parents=True, exist_ok=True)
    crop_dir.mkdir(parents=True, exist_ok=True)

    page_inputs: list[PageInput] = []
    legacy_images: list[list[tuple[str, list[float] | None]]] = []
    titles: list[str] = []

    for index, page in enumerate(document, start=1):
        rect = page.rect
        w, h = float(rect.width), float(rect.height)

        native_blocks: list[tuple[str, list[float] | None, str]] = []
        for b in page.get_text("blocks"):
            # (x0, y0, x1, y1, text, block_no, block_type); block_type 0 == text
            if len(b) >= 7 and b[6] == 0:
                text = b[4].strip()
                if text:
                    native_blocks.append((
                        text,
                        [round(b[0] / w, 4), round(b[1] / h, 4),
                         round(b[2] / w, 4), round(b[3] / h, 4)],
                        "body",
                    ))

        # Extract embedded images as visual assets.
        embedded: list[tuple[str, list[float] | None]] = []
        has_full_page_image = False
        for k, img in enumerate(page.get_images(full=True), start=1):
            xref = img[0]
            try:
                ibox = page.get_image_bbox(img)
                norm_bbox = [round(ibox.x0 / w, 4), round(ibox.y0 / h, 4),
                             round(ibox.x1 / w, 4), round(ibox.y1 / h, 4)]
            except Exception:
                norm_bbox = None
            if norm_bbox and _bbox_area(norm_bbox) >= 0.8:
                # A page-sized raster is the scanned page itself, not a figure.
                has_full_page_image = True
                continue
            try:
                extracted = document.extract_image(xref)
                ext = (extracted.get("ext") or "png").lower()
                filename = f"page_{index}_img_{k}.{ext}"
                (crop_dir / filename).write_bytes(extracted["image"])
                embedded.append((f"assets/ocr_crops/{filename}", norm_bbox))
            except Exception:
                continue

        use_native = _native_text_is_reliable(native_blocks, has_full_page_image) and force_engine is None
        render_path: Path | None = None
        input_width, input_height = w, h
        if not use_native:
            # Only rasterize pages that actually need OCR. 3x (~216 dpi) keeps
            # small Chinese text legible without multiplying every digital PDF.
            pix = page.get_pixmap(matrix=fitz.Matrix(3.0, 3.0))
            render_path = render_dir / f"page_{index}.png"
            pix.save(str(render_path))
            input_width, input_height = pix.width, pix.height

        page_inputs.append(PageInput(
            page_number=index,
            width=input_width,
            height=input_height,
            native_blocks=native_blocks if use_native else [],
            image=str(render_path) if render_path else None,
            embedded_images=embedded,
            has_native_text=use_native,
            is_whole_image=False,
        ))
        legacy_images.append(embedded)
        titles.append(native_blocks[0][0].splitlines()[0] if use_native else f"PDF 第 {index} 页")

    return _assemble(page_inputs, legacy_images, titles, "pdf", original_filename, project_root, force_engine)


# ── Scanned page image ────────────────────────────────────────────────────────

def parse_image(
    file_path: Path,
    project_root: Path,
    original_filename: str,
    force_engine: str | None = None,
) -> SourceMaterial:
    try:
        import fitz
    except ImportError as exc:
        raise ValueError("PyMuPDF is required for image parsing") from exc

    render_dir = project_root / "assets" / "ocr_render"
    render_dir.mkdir(parents=True, exist_ok=True)

    document = fitz.open(str(file_path))
    page = document[0]
    rect = page.rect
    pix = page.get_pixmap(matrix=fitz.Matrix(3.0, 3.0))
    render_path = render_dir / "page_1.png"
    pix.save(str(render_path))

    page_input = PageInput(
        page_number=1,
        width=pix.width,
        height=pix.height,
        native_blocks=[],
        image=str(render_path),
        embedded_images=[],
        has_native_text=False,
        is_whole_image=True,
    )
    return _assemble([page_input], [[]], [original_filename], "image", original_filename, project_root, force_engine)


# ── Assembly ──────────────────────────────────────────────────────────────────

def _assemble(
    page_inputs: list[PageInput],
    legacy_images: list[list[tuple[str, list[float] | None]]],
    titles: list[str],
    source_type: str,
    original_filename: str,
    project_root: Path,
    force_engine: str | None = None,
) -> SourceMaterial:
    policy = OCRPolicy(preferred_ocr=force_engine) if force_engine else None  # type: ignore[arg-type]
    analysis = run_source_understanding(page_inputs, source_type, project_root, policy=policy)

    pages: list[SourcePage] = []
    for page_input, pr, embedded, title in zip(page_inputs, analysis.pages, legacy_images, titles):
        recognized_title = next(
            (block.text for block in pr.blocks if block.block_type in TITLE_BLOCK_TYPES and block.text.strip()),
            "",
        )
        images = [
            ImageBlock(id=f"src_p{page_input.page_number}_{k}", path=rel, filename=Path(rel).name)
            for k, (rel, _) in enumerate(embedded, start=1)
        ]
        text_blocks = [
            TextBlock(
                id=b.id,
                text=b.text,
                kind="title" if b.block_type in TITLE_BLOCK_TYPES else "body",
            )
            for b in sorted(pr.blocks, key=lambda x: x.reading_order or 0)
            if b.text.strip()
        ]
        pages.append(SourcePage(
            page_number=page_input.page_number,
            title=recognized_title or title or f"第 {page_input.page_number} 页",
            text_blocks=text_blocks,
            images=images,
            ocr_text=_derive_ocr_text(pr.blocks),
        ))

    return SourceMaterial(
        source_type=source_type,  # type: ignore[arg-type]
        original_filename=original_filename,
        pages=pages,
        source_analysis=analysis,
    )


def _derive_ocr_text(blocks) -> str:
    return "\n".join(
        b.text for b in sorted(blocks, key=lambda x: x.reading_order or 0) if b.text.strip()
    )


def _bbox_area(bbox: list[float]) -> float:
    return max(0.0, bbox[2] - bbox[0]) * max(0.0, bbox[3] - bbox[1])


def _native_text_is_reliable(
    blocks: list[tuple[str, list[float] | None, str]],
    has_full_page_image: bool = False,
) -> bool:
    """Reject scan overlays and obviously corrupt PDF text layers.

    A page-sized raster is treated as a scan even if it carries hidden OCR
    text. For ordinary digital PDFs, reject control/replacement-heavy output
    and text with almost no letters or numbers.
    """
    if not blocks or has_full_page_image:
        return False
    text = "".join(text for text, _, _ in blocks if text).strip()
    if not text:
        return False
    compact = [ch for ch in text if not ch.isspace()]
    if not compact:
        return False
    invalid = sum(ch == "\ufffd" or unicodedata.category(ch) == "Cc" for ch in compact)
    signal = sum(ch.isalnum() for ch in compact)
    return invalid / len(compact) <= 0.01 and signal / len(compact) >= 0.3


def _image_size(path: Path) -> tuple[float, float] | None:
    try:
        import fitz

        doc = fitz.open(path)
        rect = doc[0].rect
        return float(rect.width), float(rect.height)
    except Exception:
        return None


def _shape_bbox_norm(shape, slide_w: float, slide_h: float) -> list[float] | None:
    left = _emu_to_points(getattr(shape, "left", None))
    top = _emu_to_points(getattr(shape, "top", None))
    width = _emu_to_points(getattr(shape, "width", None))
    height = _emu_to_points(getattr(shape, "height", None))
    if None in (left, top, width, height):
        return None
    return [
        round(left / slide_w, 4), round(top / slide_h, 4),
        round((left + width) / slide_w, 4), round((top + height) / slide_h, 4),
    ]


def _emu_to_points(value: object) -> float | None:
    if value is None:
        return None
    try:
        return round(float(value) / 12700, 2)
    except (TypeError, ValueError):
        return None
